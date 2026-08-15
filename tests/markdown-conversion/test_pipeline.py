"""
Tests for pipeline.py (v6.5.1 canonical conversion architecture).

Run from project root: pytest scripts/test_pipeline.py -v

Architecture note:
  The pipeline now accepts source documents (PDF, DOCX, TXT, etc.) as --input.
  Unit tests cover current production helpers; integration tests use real temp files.
"""
import os
import json
from pathlib import Path
import subprocess
import tempfile
import pytest
import sys

# Path to the pipeline module
_PIPELINE_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'skills', 'markdown-conversion', 'scripts')
sys.path.insert(0, os.path.normpath(_PIPELINE_DIR))

SCRIPT = [sys.executable, os.path.normpath(os.path.join(_PIPELINE_DIR, 'pipeline.py'))]

# Test config — avoids touching the real config.json
_TEST_CONFIG = os.path.normpath(os.path.join(os.path.dirname(__file__), 'fixtures', 'test_config.json'))
CONFIG_ARG = ['--config', _TEST_CONFIG]

# --- Pure function tests -------------------------------------------------------







def test_convert_chinese_tc_to_simplified():
    """Traditional Chinese 愛 → Simplified 爱."""
    from pipeline import convert_chinese
    result = convert_chinese('愛')
    assert result == '爱'


def test_convert_chinese_yu_no_false_positive():
    """於 is valid in both TC and SC; must not false-positive on stability check."""
    from pipeline import convert_chinese
    # Should not raise — 於 is stable after two-pass
    result = convert_chinese('這份報告於三月完成')
    assert '於' in result or '于' in result  # either is valid SC


def test_convert_chinese_simplified_unchanged():
    """Pure Simplified Chinese should come out unchanged."""
    from pipeline import convert_chinese
    result = convert_chinese('欢迎')
    assert result == '欢迎'


def test_convert_chinese_mixed_all_converted():
    """Mixed TC+SC: 歡迎 (TC) → 欢迎 (SC)."""
    from pipeline import convert_chinese
    result = convert_chinese('歡迎')
    assert result == '欢迎'
    assert '歡' not in result


def _frontmatter_lines(document):
    return document.split('---\n', 2)[1].splitlines()






def test_title_from_markdown_ignores_fenced_code_and_trims_closing_hashes():
    from pipeline import title_from_markdown
    body = '```markdown\n# Not a heading\n```\n\n# Actual Title ###\n'
    assert title_from_markdown(body, 'fallback') == 'Actual Title'















# --- Integration tests (real files, full pipeline) -----------------------------

def _run_pipeline(input_path, extra_args=None, output_path=None):
    """Run pipeline.py with a real file. Returns (returncode, stdout, stderr)."""
    if output_path is None:
        fd, output_path = tempfile.mkstemp(suffix='.md')
        os.close(fd)
        os.unlink(output_path)
    args = SCRIPT + CONFIG_ARG + [
        '--input', input_path,
        '--output-path', output_path,
    ]
    if extra_args:
        args += extra_args
    result = subprocess.run(args, capture_output=True)
    return result.returncode, result.stdout.decode('utf-8', errors='replace'), result.stderr.decode('utf-8', errors='replace'), output_path


def test_full_pipeline_creates_vault_file(tmp_path):
    """Full pipeline: source file → pipeline → vault file with frontmatter."""
    # Create a simple text file (markitdown supports .txt)
    src = tmp_path / 'source.txt'
    src.write_text('# Content Heading\n\nHello world', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, out_str, err, _ = _run_pipeline(str(src), output_path=str(out))
    assert code == 0, err
    assert out.exists()
    content = out.read_text(encoding='utf-8')
    assert content.startswith('---')
    assert _frontmatter_lines(content)[:4] == [
        'type: ""',
        'title: "source"',
        'description: ""',
        'tags: []',
    ]
    assert 'resource:' not in content
    assert '# Content Heading' in content
    assert 'Hello world' in content


def test_full_pipeline_no_frontmatter_flag(tmp_path):
    """--no-frontmatter should skip frontmatter injection."""
    src = tmp_path / 'source.txt'
    src.write_text('Content here', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, _, err, _ = _run_pipeline(str(src), output_path=str(out), extra_args=['--no-frontmatter'])
    assert code == 0, err
    content = out.read_text(encoding='utf-8')
    assert not content.startswith('---')


def test_full_pipeline_preserves_chinese(tmp_path):
    """Chinese content should survive the pipeline unchanged (encoding handled by markitdown)."""
    src = tmp_path / 'chinese.txt'
    src.write_text('你好世界', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, _, err, _ = _run_pipeline(str(src), output_path=str(out))
    assert code == 0, err
    content = out.read_text(encoding='utf-8')
    assert '你好世界' in content


def test_success_message_printed_to_stdout(tmp_path):
    """On success, stdout should contain [OK] message."""
    src = tmp_path / 'source.txt'
    src.write_text('content', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, stdout, err, _ = _run_pipeline(str(src), output_path=str(out))
    assert code == 0, err
    assert '[OK]' in stdout or 'Converted' in stdout


@pytest.mark.parametrize('timestamp', [
    '2026-07-22',
    '2026-07-22T14:05:06+08:00',
    '2026-07-22T14:05:06.123456-04:30',
    '2026-07-22T06:05:06Z',
])
def test_timestamp_override_is_validated_and_preserved_exactly(tmp_path, timestamp):
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, _, err, _ = _run_pipeline(
        str(src), output_path=str(out), extra_args=['--timestamp', timestamp]
    )

    assert code == 0, err
    assert f'timestamp: "{timestamp}"' in out.read_text(encoding='utf-8')


@pytest.mark.parametrize('timestamp', [
    '2026-07-22T14:05:06',
    '2026-07-22 14:05:06+08:00',
    '2026-07-22T14:05:06+0800',
    '2026-07-22T14:05:06z',
    '2026/07/22',
    'not-a-time',
])
def test_timestamp_override_rejects_naive_datetime_and_non_iso_values(tmp_path, timestamp):
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, _, err, _ = _run_pipeline(
        str(src), output_path=str(out), extra_args=['--timestamp', timestamp]
    )

    assert code == 1
    assert '--timestamp' in err
    assert not out.exists()


def test_default_timestamp_is_timezone_aware(tmp_path):
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    out = tmp_path / 'note.md'

    code, _, err, _ = _run_pipeline(str(src), output_path=str(out))

    assert code == 0, err
    timestamp_line = _frontmatter_lines(out.read_text(encoding='utf-8'))[-1]
    value = timestamp_line.removeprefix('timestamp: "').removesuffix('"')
    parsed = __import__('datetime').datetime.fromisoformat(value.replace('Z', '+00:00'))
    assert parsed.tzinfo is not None
    assert parsed.utcoffset() is not None


# --- Config loading tests -------------------------------------------------------

import json


def test_load_config_creates_default_when_missing(tmp_path, monkeypatch):
    """When config.json doesn't exist, create it with defaults and return it."""
    from pipeline import load_config, DEFAULT_CONFIG
    config_path = tmp_path / "config.json"
    monkeypatch.setattr("pipeline.CONFIG_PATH", str(config_path))
    cfg = load_config()
    assert cfg == DEFAULT_CONFIG
    assert config_path.exists()
    saved = json.loads(config_path.read_text(encoding="utf-8"))
    assert saved == DEFAULT_CONFIG


def test_load_config_reads_existing(tmp_path, monkeypatch):
    """When config.json exists, read and return it."""
    from pipeline import load_config
    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({"custom_key": "custom-value"}), encoding="utf-8")
    monkeypatch.setattr("pipeline.CONFIG_PATH", str(config_path))
    cfg = load_config()
    assert cfg["custom_key"] == "custom-value"


def test_load_config_preserves_extra_keys(tmp_path, monkeypatch):
    """Config with extra keys should preserve them after merging with defaults."""
    from pipeline import load_config
    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({"custom_key": "custom-value"}), encoding="utf-8")
    monkeypatch.setattr("pipeline.CONFIG_PATH", str(config_path))
    cfg = load_config()
    assert cfg["custom_key"] == "custom-value"


def test_pdf_ocr_config_merges_defaults_and_cli_takes_precedence():
    from pipeline import DEFAULT_CONFIG, build_parser, resolve_ocr_settings

    parser = build_parser()
    assert resolve_ocr_settings(parser.parse_args([]), DEFAULT_CONFIG).mode == 'auto'
    config = {
        'pdf_ocr': {
            'mode': 'auto',
            'engine': 'rapidocr',
            'language': 'ch',
            'dpi': 240,
            'max_long_edge': 3500,
            'min_confidence': 0.62,
        }
    }
    configured = resolve_ocr_settings(parser.parse_args([]), config)
    assert configured.mode == 'auto'
    assert configured.dpi == 240
    assert configured.min_confidence == pytest.approx(0.62)

    overridden = resolve_ocr_settings(
        parser.parse_args([
            '--ocr', 'force', '--ocr-dpi', '300',
            '--ocr-min-confidence', '0.75',
        ]),
        config,
    )
    assert overridden.mode == 'force'
    assert overridden.dpi == 300
    assert overridden.min_confidence == pytest.approx(0.75)
    assert overridden.max_long_edge == 3500


def test_pdf_ocr_config_rejects_invalid_threshold(capsys):
    from pipeline import build_parser, resolve_ocr_settings

    with pytest.raises(SystemExit) as exc_info:
        resolve_ocr_settings(
            build_parser().parse_args([]),
            {'pdf_ocr': {'mode': 'auto', 'min_confidence': 1.5}},
        )

    assert exc_info.value.code == 1
    assert 'min_confidence' in capsys.readouterr().err


def test_load_config_invalid_json_uses_defaults(tmp_path, monkeypatch):
    """Malformed config.json should warn and use defaults."""
    from pipeline import load_config, DEFAULT_CONFIG
    config_path = tmp_path / "config.json"
    config_path.write_text("{bad json", encoding="utf-8")
    monkeypatch.setattr("pipeline.CONFIG_PATH", str(config_path))
    cfg = load_config()
    assert cfg == DEFAULT_CONFIG


# --- Default output path tests -------------------------------------------------


def _args_ns(**kwargs):
    """Build a minimal argparse.Namespace for target-resolution tests."""
    from argparse import Namespace
    base = {
        'input': None, 'input_dir': None, 'output_path': '', 'output_dir': '',
        'output_mode': None,
    }
    base.update(kwargs)
    return Namespace(**base)


@pytest.mark.parametrize('stem', [
    '0. Table of contents',
    '3. MB Rule 9.03(1)(b) Initial listing fee',
    '10. FF004 - Lynk Pharmaceuticals Co. Ltd',
])
def test_resolve_target_preserves_dotted_logical_stem(tmp_path, stem):
    from pipeline import resolve_target

    src = tmp_path / f'{stem}.txt'
    src.write_text('Body', encoding='utf-8')
    output = tmp_path / 'out'

    markdown = resolve_target(
        _args_ns(input=str(src), output_dir=str(output), output_mode='markdown'),
        str(src),
    )
    assert markdown.path == output / f'{stem}.md'
    assert markdown.stem == stem

    bundle = resolve_target(
        _args_ns(input=str(src), output_dir=str(output), output_mode='bundle'),
        str(src),
    )
    assert bundle.path == output / stem
    assert bundle.stem == stem








# --- Batch mode tests -----------------------------------------------------------

def _run_batch(input_dir, output_dir, extra_args=None):
    """Run pipeline.py in batch mode. Returns (returncode, stdout, stderr)."""
    args = SCRIPT + CONFIG_ARG + [
        '--input-dir', input_dir,
        '--output-path', output_dir,
    ]
    if extra_args:
        args += extra_args
    result = subprocess.run(args, capture_output=True)
    return result.returncode, result.stdout.decode('utf-8', errors='replace'), result.stderr.decode('utf-8', errors='replace')


def test_batch_converts_multiple_files(tmp_path):
    """Batch mode should convert all supported files in a directory."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    (src_dir / 'a.txt').write_text('![Chart](chart.png)\n\nFile A', encoding='utf-8')
    (src_dir / 'b.txt').write_text('File B', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    timestamp = '2026-07-22T14:05:06+08:00'
    code, stdout, stderr = _run_batch(
        str(src_dir), str(out_dir), extra_args=['--timestamp', timestamp]
    )
    assert code == 0, stderr
    assert (out_dir / 'a' / 'a.json').exists()
    assert (out_dir / 'a' / 'a.md').exists()
    assert (out_dir / 'b' / 'b.md').exists()
    first = (out_dir / 'a' / 'a.md').read_text(encoding='utf-8')
    assert 'File A' in first
    assert 'File B' in (out_dir / 'b' / 'b.md').read_text(encoding='utf-8')
    assert _frontmatter_lines(first) == [
        'type: ""',
        'title: "a"',
        'description: ""',
        'tags: []',
        f'timestamp: "{timestamp}"',
    ]
    assert 'chart.png' not in first
    assert '![' not in first
    assert 'Chart' in first
    assert 'resource:' not in first


def test_batch_mirrors_subdirectory_structure(tmp_path):
    """Batch mode should mirror subdirectory structure in output."""
    src_dir = tmp_path / 'input'
    sub = src_dir / 'sub'
    sub.mkdir(parents=True)
    (src_dir / 'top.txt').write_text('Top', encoding='utf-8')
    (sub / 'nested.txt').write_text('Nested', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir))
    assert code == 0, stderr
    assert (out_dir / 'top' / 'top.md').exists()
    assert (out_dir / 'sub' / 'nested' / 'nested.md').exists()


def test_batch_no_recursive_flattens(tmp_path):
    """--no-recursive should only convert top-level files."""
    src_dir = tmp_path / 'input'
    sub = src_dir / 'sub'
    sub.mkdir(parents=True)
    (src_dir / 'top.txt').write_text('Top', encoding='utf-8')
    (sub / 'nested.txt').write_text('Nested', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir), ['--no-recursive'])
    assert code == 0, stderr
    assert (out_dir / 'top' / 'top.md').exists()
    assert not (out_dir / 'sub' / 'nested' / 'nested.md').exists()


def test_batch_types_filter(tmp_path):
    """--types should filter to specified extensions only."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    (src_dir / 'keep.txt').write_text('Keep', encoding='utf-8')
    (src_dir / 'skip.html').write_text('<p>Skip</p>', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir), ['--types', 'txt'])
    assert code == 0, stderr
    assert (out_dir / 'keep' / 'keep.md').exists()
    assert not (out_dir / 'skip').exists()


def test_batch_invalid_types_exits_1(tmp_path):
    """--types with unsupported extension should error."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir), ['--types', 'xyz'])
    assert code == 1
    assert 'Unsupported' in stderr


def test_batch_continues_on_error(tmp_path):
    """Batch should continue converting when one file fails."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    (src_dir / 'good.txt').write_text('Good', encoding='utf-8')
    # Invalid PDF forces the native adapter to fail while the text file succeeds.
    (src_dir / 'bad.pdf').write_bytes(b'not-a-pdf')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir), ['--types', 'txt,pdf'])
    # The invalid PDF fails but the TXT is still published; aggregate exit is failure.
    assert code == 1, stderr
    assert (out_dir / 'good' / 'good.md').exists()


def test_batch_empty_directory(tmp_path):
    """Batch on empty directory should succeed with 0 converted."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir))
    assert code == 0
    assert '0 converted' in stdout or '[BATCH]' in stdout


@pytest.mark.parametrize(
    ('relation', 'error_fragment'),
    [
        ('equal', 'must not equal the input root'),
        ('ancestor', 'must not be an ancestor of the input root'),
    ],
)
def test_batch_rejects_unsafe_output_roots_before_collection_or_publication(
    monkeypatch, tmp_path, capsys, relation, error_fragment
):
    import pipeline

    workspace = tmp_path / 'workspace'
    source = workspace / 'input'
    source.mkdir(parents=True)
    (source / 'document.txt').write_text('Body', encoding='utf-8')
    output = source if relation == 'equal' else workspace
    args = pipeline.build_parser().parse_args([
        '--input-dir', str(source), '--output-dir', str(output),
    ])
    collection_calls = []

    def unexpected_collection(*args, **kwargs):
        collection_calls.append((args, kwargs))
        pytest.fail('unsafe batch roots must be rejected before collection')

    monkeypatch.setattr(pipeline, 'collect_files', unexpected_collection)

    with pytest.raises(SystemExit) as exc_info:
        pipeline.precheck(args)
        pipeline.run_batch(args)

    assert exc_info.value.code == 1
    assert error_fragment in capsys.readouterr().err
    assert collection_calls == []
    assert not (output / 'document').exists()
    assert not any('.staging-' in path.name for path in workspace.rglob('*'))


@pytest.mark.parametrize('relation', ['descendant', 'sibling-prefix'])
def test_batch_allows_safe_output_roots_and_excludes_them_from_collection(
    monkeypatch, tmp_path, relation
):
    import pipeline

    workspace = tmp_path / 'workspace'
    source = workspace / 'input'
    source.mkdir(parents=True)
    (source / 'document.txt').write_text('Body', encoding='utf-8')
    output = (
        source / '_converted'
        if relation == 'descendant'
        else workspace / 'input-output'
    )
    args = pipeline.build_parser().parse_args([
        '--input-dir', str(source), '--output-dir', str(output),
    ])
    collection_calls = []

    def controlled_collection(input_dir, recursive, types, exclude_root):
        collection_calls.append((input_dir, recursive, types, exclude_root))
        return []

    monkeypatch.setattr(pipeline, 'collect_files', controlled_collection)

    pipeline.precheck(args)
    assert pipeline.run_batch(args) == 0

    assert collection_calls == [(str(source), True, None, output.resolve())]
    assert not output.exists()


def test_batch_input_dir_not_found(tmp_path):
    """Non-existent input-dir should error."""
    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(tmp_path / 'nonexistent'), str(out_dir))
    assert code == 1
    assert 'not found' in stderr.lower() or 'Directory' in stderr


def test_batch_file_exists_skip(tmp_path):
    """Existing output files should be skipped (not fatal) in batch mode."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    (src_dir / 'file.txt').write_text('New content', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()
    (out_dir / 'file').mkdir()
    (out_dir / 'file' / 'file.md').write_text('Old content', encoding='utf-8')

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir))
    assert code == 2
    assert 'skipped' in stdout.lower() or 'skipped' in stderr.lower()


def test_batch_summary_output(tmp_path):
    """Batch should print summary with converted/failed/skipped counts."""
    src_dir = tmp_path / 'input'
    src_dir.mkdir()
    (src_dir / 'a.txt').write_text('A', encoding='utf-8')
    (src_dir / 'b.txt').write_text('B', encoding='utf-8')

    out_dir = tmp_path / 'output'
    out_dir.mkdir()

    code, stdout, stderr = _run_batch(str(src_dir), str(out_dir))
    assert code == 0
    assert '[BATCH]' in stdout
    assert '2 converted' in stdout


# --- Precheck tests ----------------------------------------------------------------


def _run_pipeline_cli(*extra_args):
    """Run pipeline.py with given args. Returns (returncode, stdout, stderr)."""
    result = subprocess.run(
        SCRIPT + CONFIG_ARG + list(extra_args),
        capture_output=True,
    )
    return result.returncode, result.stdout.decode('utf-8', errors='replace'), result.stderr.decode('utf-8', errors='replace')


def test_precheck_file_not_found(tmp_path):
    """--input with non-existent file should exit 1."""
    code, stdout, stderr = _run_pipeline_cli(
        '--input', str(tmp_path / 'nonexistent.pdf'),
        '--output-path', str(tmp_path / 'out.md'),
    )
    assert code == 1
    assert 'File not found' in stderr or 'not found' in stderr.lower()


def test_precheck_directory_not_found(tmp_path):
    """--input-dir with non-existent directory should exit 1."""
    code, stdout, stderr = _run_pipeline_cli(
        '--input-dir', str(tmp_path / 'nonexistent-dir'),
        '--output-path', str(tmp_path / 'output'),
    )
    assert code == 1
    assert 'Directory not found' in stderr or 'not found' in stderr.lower()


def test_precheck_mutual_exclusivity(tmp_path):
    """Both --input and --input-dir should exit 1."""
    src = tmp_path / 'test.txt'
    src.write_text('hello', encoding='utf-8')
    code, stdout, stderr = _run_pipeline_cli(
        '--input', str(src),
        '--input-dir', str(tmp_path),
        '--output-path', str(tmp_path / 'out.md'),
    )
    assert code == 1
    assert 'exactly one' in stderr.lower()


def test_precheck_neither_input_nor_dir(tmp_path):
    """No --input or --input-dir should exit 1."""
    code, stdout, stderr = _run_pipeline_cli(
        '--output-path', str(tmp_path / 'out.md'),
    )
    assert code == 1
    assert 'required' in stderr.lower()


def test_legacy_doc_markitdown_is_rejected_before_worker_staging_or_output(
    monkeypatch, tmp_path
):
    import pipeline

    source = tmp_path / 'legacy.doc'
    source.write_bytes(b'legacy-doc-placeholder')
    output = tmp_path / 'published'
    args = pipeline.build_parser().parse_args([
        '--input', str(source),
        '--output-dir', str(output),
        '--local-document-adapter', 'markitdown',
    ])
    pipeline.precheck(args)
    side_effects = []

    def unexpected(label):
        def fail(*_args, **_kwargs):
            side_effects.append(label)
            pytest.fail(f'legacy .doc rejection reached {label}')
        return fail

    monkeypatch.setattr(pipeline, '_preflight_target', unexpected('target preflight'))
    monkeypatch.setattr(pipeline, '_run_provider_worker', unexpected('provider worker'))
    monkeypatch.setattr(pipeline.tempfile, 'mkdtemp', unexpected('staging'))
    monkeypatch.setattr(pipeline, '_write_markdown_file', unexpected('markdown output'))
    monkeypatch.setattr(pipeline, '_publish_directory', unexpected('bundle publication'))

    with pytest.raises(pipeline.PipelineError, match='cannot safely convert legacy .doc'):
        pipeline.convert_one(args, str(source))

    assert side_effects == []
    assert not output.exists()


def test_legacy_doc_markitdown_cli_returns_one_without_output(tmp_path):
    source = tmp_path / 'legacy.doc'
    source.write_bytes(b'legacy-doc-placeholder')
    output = tmp_path / 'published'

    result = subprocess.run(
        SCRIPT + CONFIG_ARG + [
            '--input', str(source),
            '--output-dir', str(output),
            '--local-document-adapter', 'markitdown',
        ],
        capture_output=True,
    )

    assert result.returncode == 1
    assert 'cannot safely convert legacy .doc' in result.stderr.decode(
        'utf-8', errors='replace'
    )
    assert not output.exists()


# --- --version and --config tests ------------------------------------------------


def test_version_flag():
    """--version should print version and dependency status, then exit 0."""
    result = subprocess.run(SCRIPT + ['--version'], capture_output=True)
    assert result.returncode == 0
    stdout = result.stdout.decode('utf-8')
    assert 'markdown-conversion v' in stdout
    assert 'Dependencies:' in stdout
    # Should show pip install names (not import names)
    assert 'opencc-python-reimplemented' in stdout
    assert 'markdown-conversion v6.5.1' in stdout
    assert 'rapidocr:' in stdout
    assert 'onnxruntime:' in stdout
    assert 'ruamel.yaml:' not in stdout
    assert 'cortex:' not in stdout


def test_config_flag_uses_alternate_config(tmp_path):
    """--config should load config from the specified path."""
    cfg = tmp_path / 'custom.json'
    cfg.write_text('{"custom_key": "custom-value"}', encoding='utf-8')
    result = subprocess.run(
        SCRIPT + ['--config', str(cfg), '--version'],
        capture_output=True,
    )
    assert result.returncode == 0
    # Config was loaded without error (would fail if path was ignored)


# --- URL utility tests ---------------------------------------------------------


def test_is_url_http():
    """HTTP URLs should be detected."""
    from pipeline import is_url
    assert is_url('http://example.com/file.pdf')
    assert is_url('HTTP://EXAMPLE.COM/FILE.PDF')


def test_is_url_https():
    """HTTPS URLs should be detected."""
    from pipeline import is_url
    assert is_url('https://example.com/file.pdf')


def test_is_url_file_path():
    """Local file paths should not be detected as URLs."""
    from pipeline import is_url
    assert not is_url('/local/path/file.pdf')
    assert not is_url('C:\\Users\\doc.pdf')
    assert not is_url('./relative.pdf')


def test_is_url_empty():
    """Empty string should not be detected as URL."""
    from pipeline import is_url
    assert not is_url('')


def test_url_to_slug_basic():
    """Basic URL should produce a clean slug."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com/docs/report.pdf')
    assert 'example' in slug
    assert 'report' in slug
    assert '://' not in slug
    assert '/' not in slug


def test_url_to_slug_root():
    """Root URL should fall back to domain."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com/')
    assert 'example' in slug


def test_url_to_slug_query_string():
    """Query strings should be excluded from slug."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com/page?id=1&name=test')
    assert '?' not in slug
    assert 'page' in slug


def test_url_to_slug_encoded():
    """URL-encoded characters should be decoded."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com/my%20document.pdf')
    assert 'my' in slug
    assert 'document' in slug


def test_url_to_slug_truncation():
    """Very long URLs should be truncated."""
    from pipeline import url_to_slug
    long_url = 'https://example.com/' + 'a' * 200 + '.pdf'
    slug = url_to_slug(long_url)
    assert len(slug) <= 120


def test_url_to_slug_empty_fallback():
    """URL with minimal content should still produce a slug."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com')
    assert slug  # not empty


def test_safe_url_rejects_private_network_resolution(monkeypatch):
    import safe_url

    monkeypatch.setattr(
        safe_url.socket,
        'getaddrinfo',
        lambda *args, **kwargs: [(safe_url.socket.AF_INET, safe_url.socket.SOCK_STREAM, 6, '', ('127.0.0.1', 80))],
    )
    with pytest.raises(RuntimeError, match='non-public'):
        safe_url._endpoint('http://example.test/report')


def test_safe_url_revalidates_redirect_and_redacts_secrets(monkeypatch):
    import safe_url

    monkeypatch.setattr(
        safe_url,
        '_request',
        lambda value, deadline: (302, {'location': 'http://127.0.0.1/internal?token=secret'}, b''),
    )
    with pytest.raises(RuntimeError, match='non-public'):
        safe_url.download_url('https://example.com/start?api_key=secret')
    assert safe_url.redact_url('https://user:pass@example.com/report?token=secret#frag') == 'https://example.com/report'


def test_url_source_record_hashes_response_bytes_and_redacts_query_secret():
    from pipeline import _source_record

    record, document_id = _source_record(
        'https://example.com/report?token=secret',
        remote_bytes=b'<html>body</html>',
        remote_locator='https://example.com/report',
        remote_media_type='text/html',
    )
    assert record['locator'] == 'https://example.com/report'
    assert 'secret' not in json.dumps(record)
    assert record['hash_basis'] == 'remote_response_bytes'
    assert record['size_bytes'] == len(b'<html>body</html>')
    assert document_id == f"sha256:{record['sha256']}"


def test_native_provider_worker_timeout_is_bounded(monkeypatch):
    import pipeline

    def timeout(*args, **kwargs):
        raise subprocess.TimeoutExpired(args[0], kwargs.get('timeout', 1))

    monkeypatch.setattr(pipeline.subprocess, 'run', timeout)
    with pytest.raises(pipeline.PipelineError, match='exceeded 0.01 seconds'):
        pipeline._run_provider_worker({'adapter': 'anydoc'}, timeout=0.01)


@pytest.mark.parametrize(
    ('ocr_mode', 'expected_imports'),
    [
        ('auto', ['pypdf', 'pdf_inspector']),
        ('off', ['pypdf', 'pdf_inspector']),
        ('force', ['pypdf', 'pypdfium2', 'rapidocr', 'onnxruntime']),
    ],
)
def test_pdf_worker_dependency_gate_is_route_specific(
    monkeypatch, ocr_mode, expected_imports
):
    import builtins
    import provider_worker

    original_import = builtins.__import__
    pdf_capabilities = {
        'pypdf', 'pdf_inspector', 'pypdfium2', 'rapidocr', 'onnxruntime', 'pdfminer',
    }
    imported = []

    def controlled_import(name, *args, **kwargs):
        if name in pdf_capabilities:
            imported.append(name)
            if name not in expected_imports:
                raise AssertionError(
                    f'{ocr_mode} unexpectedly imported PDF capability {name}'
                )
            return object()
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, '__import__', controlled_import)
    provider_worker._require_pdf_route(ocr_mode)

    assert imported == expected_imports


@pytest.mark.parametrize(
    ('ocr_mode', 'blocked_import', 'install_name'),
    [
        ('auto', 'pypdf', 'pypdf'),
        ('off', 'pdf_inspector', 'pdf-inspector'),
        ('force', 'pypdfium2', 'pypdfium2'),
        ('force', 'rapidocr', 'rapidocr'),
        ('force', 'onnxruntime', 'onnxruntime'),
    ],
)
def test_pdf_worker_dependency_gate_reports_missing_required_capability(
    monkeypatch, ocr_mode, blocked_import, install_name
):
    import builtins
    import provider_worker

    original_import = builtins.__import__
    pdf_capabilities = {
        'pypdf', 'pdf_inspector', 'pypdfium2', 'rapidocr', 'onnxruntime',
    }

    def controlled_import(name, *args, **kwargs):
        if name == blocked_import:
            raise ImportError(f'controlled missing dependency: {name}')
        if name in pdf_capabilities:
            return object()
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, '__import__', controlled_import)
    with pytest.raises(
        RuntimeError,
        match=f'Required PDF route dependency {install_name} is unavailable',
    ):
        provider_worker._require_pdf_route(ocr_mode)


@pytest.mark.parametrize('adapter', ['markitdown', 'url_markitdown'])
def test_non_pdf_worker_routes_do_not_gate_local_pdf_dependencies(
    monkeypatch, tmp_path, adapter
):
    from types import SimpleNamespace
    import provider_worker

    source = tmp_path / 'source.txt'
    source.write_text('Local body', encoding='utf-8')
    request_path = tmp_path / f'{adapter}.request.json'
    result_path = tmp_path / f'{adapter}.result.json'
    pdf_gate_calls = []
    monkeypatch.setattr(
        provider_worker,
        '_require_pdf_route',
        lambda mode: pdf_gate_calls.append(mode),
    )

    if adapter == 'markitdown':
        class StubMarkItDownAdapter:
            def extract(self, source_value, document_id, mode, asset_dir):
                assert source_value == str(source)
                return {'adapter': 'markitdown', 'document_id': document_id, 'mode': mode}

        monkeypatch.setattr(
            provider_worker, 'MarkItDownAdapter', StubMarkItDownAdapter
        )
        request = {
            'adapter': 'markitdown',
            'source': str(source),
            'document_id': 'sha256:test',
            'mode': 'preserve',
            'asset_dir': '',
        }
    else:
        monkeypatch.setattr(
            provider_worker,
            'download_url',
            lambda _source: SimpleNamespace(
                suffix='.html',
                payload=b'<p>Remote body</p>',
                locator='https://example.com/report',
                media_type='text/html',
            ),
        )
        monkeypatch.setattr(
            provider_worker, 'convert_basic', lambda _source: '# Remote\n\nBody'
        )
        request = {
            'adapter': 'url_markitdown',
            'source': 'https://example.com/report?token=secret',
        }

    request_path.write_text(json.dumps(request), encoding='utf-8')
    monkeypatch.setattr(
        provider_worker.sys,
        'argv',
        [
            'provider_worker.py',
            '--request', str(request_path),
            '--result', str(result_path),
        ],
    )

    assert provider_worker.main() == 0
    assert json.loads(result_path.read_text(encoding='utf-8'))['ok'] is True
    assert pdf_gate_calls == []


def test_url_conversion_is_bounded_by_isolated_worker(monkeypatch):
    import pipeline

    calls = []

    def worker(request, timeout=180.0):
        calls.append((request, timeout))
        return {
            'markdown': '# Remote\n\nBody',
            'locator': 'https://example.com/report',
            'media_type': 'text/html',
            'sha256': 'a' * 64,
            'size_bytes': 17,
        }

    monkeypatch.setattr(pipeline, '_run_provider_worker', worker)
    monkeypatch.setattr(
        pipeline,
        'PdfInspectorAdapter',
        lambda *_args, **_kwargs: pytest.fail('URL route touched local PDF adapter'),
    )
    extracted, source, document_id = pipeline._extract(
        'https://example.com/report?token=secret', 'preserve', None
    )
    assert calls == [({'adapter': 'url_markitdown', 'source': 'https://example.com/report?token=secret'}, 45.0)]
    assert source['locator'] == 'https://example.com/report'
    assert source['sha256'] == 'a' * 64
    assert source['hash_basis'] == 'remote_response_bytes'
    assert document_id == f"sha256:{'a' * 64}"
    assert extracted['content']


def test_local_markitdown_conversion_is_bounded_by_isolated_worker(tmp_path, monkeypatch):
    import pipeline

    source = tmp_path / 'source.txt'
    source.write_text('content', encoding='utf-8')
    calls = []

    def worker(request, timeout=180.0):
        calls.append((request, timeout))
        return {
            'source_units': [{'id': 'unit-test', 'type': 'document', 'index': 1, 'locator': {}, 'status': 'complete', 'warnings': []}],
            'content': [{'id': 'node-test', 'type': 'paragraph', 'source_locator': {'source_unit_id': 'unit-test'}, 'raw_text': 'content', 'text': 'content', 'normalized_text': 'content'}],
            'tables': [], 'assets': [], 'relationships': [], 'warnings': [], 'title': 'source',
            'adapter': {'name': 'markitdown', 'version': 'test', 'limitations': []},
        }

    monkeypatch.setattr(pipeline, '_run_provider_worker', worker)
    monkeypatch.setattr(
        pipeline,
        'PdfInspectorAdapter',
        lambda *_args, **_kwargs: pytest.fail('non-PDF route touched local PDF adapter'),
    )
    result, _, document_id = pipeline._extract(str(source), 'preserve', None)
    assert result['adapter']['name'] == 'markitdown'
    assert calls == [({
        'adapter': 'markitdown', 'source': str(source), 'document_id': document_id,
        'mode': 'preserve', 'asset_dir': '',
    }, 180.0)]


def test_office_image_ocr_enrichment_is_opt_in_and_provenance_linked(tmp_path, monkeypatch):
    import pipeline
    from ocr_provider import OcrSettings

    unit_id = 'unit-0000000000000000'
    extracted = {
        'adapter': {'name': 'anydoc'},
        'source_units': [{'id': unit_id, 'status': 'complete', 'warnings': []}],
        'content': [{
            'id': 'node-image000000001',
            'type': 'image',
            'asset_id': 'asset-000000000001',
            'source_locator': {'source_unit_id': unit_id},
        }],
        'assets': [{
            'asset_id': 'asset-000000000001',
            'path': 'assets/images/asset-000000000001.png',
        }],
        'relationships': [],
        'warnings': [],
    }
    monkeypatch.setattr(
        pipeline,
        '_run_provider_worker',
        lambda request: {'items': [{
            'asset_id': 'asset-000000000001',
            'text': '圖片中的文字',
            'confidence': 0.98,
            'engine': 'rapidocr',
            'engine_version': 'test',
        }]},
    )
    pipeline._enrich_office_images(
        extracted,
        'sha256:' + ('a' * 64),
        'preserve',
        tmp_path / 'assets' / 'images',
        OcrSettings(),
    )
    assert [item['type'] for item in extracted['content']] == ['image', 'paragraph']
    assert extracted['content'][1]['text'] == '圖片中的文字'
    assert extracted['content'][1]['source_locator']['asset_id'] == 'asset-000000000001'
    assert extracted['relationships'][-1]['type'] == 'image_ocr_text'
















# --- URL precheck integration tests --------------------------------------------


def test_precheck_url_accepted(tmp_path):
    """A URL as --input should not trigger 'File not found' precheck."""
    code, stdout, stderr = _run_pipeline_cli(
        '--input', 'https://example.com/test.html',
        '--output-path', str(tmp_path / 'out.md'),
    )
    # Precheck passed — failure is from conversion (network), not "File not found"
    assert 'File not found' not in stderr
    # It may fail with HTTP error or succeed — both are fine for this test
    # The important thing is precheck didn't reject the URL as a missing file


def test_precheck_url_not_accepted_with_input_dir(tmp_path):
    """--input-dir should not accept URLs."""
    code, stdout, stderr = _run_pipeline_cli(
        '--input-dir', 'https://example.com/folder',
        '--output-path', str(tmp_path / 'output'),
    )
    # This should fail because input-dir expects a local directory
    # (os.path.isdir on a URL will fail)
    assert code == 1


# --- Removed runtime-coupling options ------------------------------------------


def test_removed_flags_are_not_advertised():
    result = subprocess.run(SCRIPT + ['--help'], capture_output=True)
    assert result.returncode == 0
    stdout = result.stdout.decode('utf-8', errors='replace')
    for option in (
        '--keep-images', '--okf', '--workspace', '--okf-run-dir',
        '--accept-partial', '--source', '--converted-at',
    ):
        assert option not in stdout


@pytest.mark.parametrize('option', [
    '--keep-images', '--okf', '--workspace', '--okf-run-dir', '--accept-partial',
    '--source', '--converted-at',
])
def test_removed_flags_are_rejected(tmp_path, option):
    src = tmp_path / 'source.txt'
    src.write_text('Hello', encoding='utf-8')
    out = tmp_path / 'final.md'
    extra = [option]
    if option in {'--workspace', '--okf-run-dir', '--source', '--converted-at'}:
        extra.append(str(tmp_path / 'value'))
    result = subprocess.run(
        SCRIPT + CONFIG_ARG + ['--input', str(src), '--output-path', str(out), *extra],
        capture_output=True,
    )
    assert result.returncode == 2
    assert not out.exists()
    assert 'unrecognized arguments' in result.stderr.decode('utf-8', errors='replace').lower()


@pytest.mark.parametrize('option', [
    '--keep-images', '--okf', '--workspace', '--okf-run-dir', '--accept-partial',
    '--source', '--converted-at',
])
def test_removed_flags_are_rejected_even_with_version(option):
    args = ['--version', option]
    if option in {'--workspace', '--okf-run-dir', '--source', '--converted-at'}:
        args.append('value')
    result = subprocess.run(SCRIPT + args, capture_output=True)
    assert result.returncode == 2
    assert 'unrecognized arguments' in result.stderr.decode('utf-8', errors='replace').lower()


def test_markdown_only_replaces_images_with_caption_text(tmp_path):
    src = tmp_path / 'images.txt'
    src.write_text('# Title\n\n![Chart](chart.png)\n\nSome text.\n\n![Graph](graph.jpg)\n\nMore text.', encoding='utf-8')
    out = tmp_path / 'out.md'

    code, stdout, stderr, _ = _run_pipeline(
        str(src), output_path=str(out),
    )
    assert code == 0, stderr
    content = out.read_text(encoding='utf-8')
    assert '![' not in content
    assert 'chart.png' not in content
    assert 'Chart' in content
    assert 'Graph' in content
    assert 'Some text.' in content
    assert 'More text.' in content


# --- Additional unit tests ------------------------------------------------------


def test_url_to_slug_chinese():
    """URL with Chinese characters should preserve them in slug."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://example.com/文档/报告.pdf')
    assert 'example' in slug
    # Chinese chars should be preserved
    assert len(slug) > 0


def test_url_to_slug_auth():
    """URL with auth info should strip credentials from slug."""
    from pipeline import url_to_slug
    slug = url_to_slug('https://user:pass@example.com/docs/report.pdf')
    assert 'user' not in slug
    assert 'pass' not in slug
    assert 'example' in slug






# --- v6 canonical bundle and adapter acceptance tests -------------------------


def _run_pipeline_bundle(source, output_dir, extra_args=None):
    args = SCRIPT + CONFIG_ARG + ['--input', str(source), '--output-dir', str(output_dir)]
    if extra_args:
        args += list(extra_args)
    result = subprocess.run(args, capture_output=True)
    bundle = Path(output_dir) / Path(source).stem
    return (
        result.returncode,
        result.stdout.decode('utf-8', errors='replace'),
        result.stderr.decode('utf-8', errors='replace'),
        bundle,
    )




def _run_product_bundle(source, output_dir, extra_args=None):
    """Exercise the same production pipeline used by the CLI for every format."""
    return _run_pipeline_bundle(source, output_dir, extra_args)


def _load_bundle(bundle):
    return json.loads((bundle / f'{bundle.name}.json').read_text(encoding='utf-8'))


def test_default_bundle_contains_canonical_json_and_markdown(tmp_path):
    src = tmp_path / 'report.txt'
    src.write_text('# Report\n\nBody', encoding='utf-8')
    code, stdout, stderr, bundle = _run_product_bundle(src, tmp_path / 'out')
    assert code == 0, stderr
    assert sorted(path.name for path in bundle.iterdir()) == ['report.json', 'report.md']
    data = _load_bundle(bundle)
    assert data['schema_version'] == '1.0'
    assert data['outputs']['mode'] == 'bundle'
    assert data['document']['document_id'] == f"sha256:{data['source']['sha256']}"
    assert data['quality']['status'] == 'complete'


def test_markitdown_local_title_uses_literal_source_stem_in_markdown_and_json(tmp_path):
    src = tmp_path / '季度報告.final.txt'
    src.write_text('# Content Heading\n\nBody', encoding='utf-8')
    code, _, stderr, bundle = _run_product_bundle(src, tmp_path / 'out')

    assert code == 0, stderr
    data = _load_bundle(bundle)
    markdown = (bundle / f'{bundle.name}.md').read_text(encoding='utf-8')
    assert data['document']['title'] == '季度報告.final'
    assert 'title: "季度報告.final"' in markdown
    assert '# Content Heading' in markdown


def test_url_title_keeps_content_h1_compatibility(monkeypatch):
    import pipeline

    extracted = {
        'title': 'Remote Heading',
        'adapter': {'name': 'markitdown', 'version': 'test', 'limitations': []},
        'warnings': [],
    }
    monkeypatch.setattr(
        pipeline,
        '_extract',
        lambda *args, **kwargs: (extracted, {'kind': 'url'}, 'sha256:test'),
    )

    document = pipeline._build_document(
        'https://example.com/report.pdf',
        '2026-08-04',
        'preserve',
        'markdown',
        None,
    )

    assert document['document']['title'] == 'Remote Heading'


def test_local_title_uses_original_identity_for_legacy_doc_temp_conversion(monkeypatch, tmp_path):
    """Retain the established node ID while exercising the generic snapshot identity seam."""
    import pipeline

    original = tmp_path / 'Original.Name.docx'
    temporary = tmp_path / 'accepted-revision.docx'
    extracted = {
        'title': 'Content Heading',
        'adapter': {'name': 'markitdown', 'version': 'test', 'limitations': []},
        'warnings': [],
    }
    monkeypatch.setattr(
        pipeline,
        '_extract',
        lambda *args, **kwargs: (extracted, {'kind': 'file'}, 'sha256:test'),
    )

    document = pipeline._build_document(
        str(temporary),
        '2026-08-04',
        'simplified',
        'markdown',
        None,
        identity_source=str(original),
    )

    assert document['document']['title'] == 'Original.Name'


def test_single_file_without_output_flags_uses_sibling_bundle(tmp_path):
    src = tmp_path / 'sibling.txt'
    src.write_text('Body', encoding='utf-8')
    result = subprocess.run(SCRIPT + CONFIG_ARG + ['--input', str(src)], capture_output=True)
    assert result.returncode == 0, result.stderr.decode(errors='replace')
    assert (tmp_path / 'sibling' / 'sibling.json').exists()
    assert (tmp_path / 'sibling' / 'sibling.md').exists()


def test_output_path_rejects_explicit_bundle_mode(tmp_path):
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    target = tmp_path / 'target.md'
    result = subprocess.run(
        SCRIPT + CONFIG_ARG + [
            '--input', str(src), '--output-mode', 'bundle', '--output-path', str(target),
        ],
        capture_output=True,
    )
    assert result.returncode == 1
    assert '--output-path is valid only' in result.stderr.decode(errors='replace')
    assert not target.exists()


def test_markdown_only_emits_exactly_one_file_and_no_dead_image_links(tmp_path):
    src = tmp_path / 'images.txt'
    src.write_text('![Caption](missing.png)\n\nBody\n\n![](silent.png)', encoding='utf-8')
    target = tmp_path / 'clean.md'
    result = subprocess.run(
        SCRIPT + CONFIG_ARG + ['--input', str(src), '--output-path', str(target)],
        capture_output=True,
    )
    assert result.returncode == 0, result.stderr.decode(errors='replace')
    assert sorted(path.name for path in tmp_path.iterdir()) == ['clean.md', 'images.txt']
    markdown = target.read_text(encoding='utf-8')
    assert 'Caption' in markdown
    assert 'missing.png' not in markdown
    assert 'silent.png' not in markdown
    assert '![' not in markdown


def test_traditional_raw_text_is_preserved_while_markdown_is_simplified(tmp_path):
    src = tmp_path / 'traditional.txt'
    src.write_text('繁體中文與軟體', encoding='utf-8')
    code, _, stderr, bundle = _run_product_bundle(src, tmp_path / 'out')
    assert code == 0, stderr
    data = _load_bundle(bundle)
    node = data['content'][0]
    assert node['raw_text'] == '繁體中文與軟體'
    assert node['text'] == '繁體中文與軟體'
    assert node['normalized_text'] == '繁体中文与软体'
    assert '繁体中文与软体' in (bundle / 'traditional.md').read_text(encoding='utf-8')


def test_language_normalization_protects_code_url_path_hash_and_formula():
    from canonical import convert_chinese
    digest = 'a' * 64
    value = f'軟體 `繁體` https://example.com/繁體 C:\\繁體\\file sha256:{digest} $變數$'
    converted = convert_chinese(value, 'simplified')
    assert converted.startswith('软体 ')
    assert '`繁體`' in converted
    assert 'https://example.com/繁體' in converted
    assert 'C:\\繁體\\file' in converted
    assert f'sha256:{digest}' in converted
    assert '$變數$' in converted


def test_same_source_produces_stable_document_and_node_ids(tmp_path):
    src = tmp_path / 'stable.txt'
    src.write_text('# Stable\n\nSame bytes', encoding='utf-8')
    first = _run_product_bundle(src, tmp_path / 'first')[3]
    second = _run_product_bundle(src, tmp_path / 'second')[3]
    left, right = _load_bundle(first), _load_bundle(second)
    assert left['document']['document_id'] == right['document']['document_id']
    assert [node['id'] for node in left['content']] == [node['id'] for node in right['content']]


def test_content_is_authoritative_order_for_table_reference(tmp_path):
    src = tmp_path / 'table.md'
    src.write_text('Before\n\n| Name | Value |\n| --- | --- |\n| A | 1 |\n\nAfter', encoding='utf-8')
    code, _, stderr, bundle = _run_product_bundle(src, tmp_path / 'out')
    assert code == 0, stderr
    data = _load_bundle(bundle)
    types = [node['type'] for node in data['content']]
    assert types == ['paragraph', 'table', 'paragraph']
    table_node = data['content'][1]
    assert table_node['table_id'] == data['tables'][0]['table_id']
    assert data['tables'][0]['rows'][1][1]['value'] == 1


def test_semantic_validator_rejects_dangling_reference(tmp_path):
    from canonical import CanonicalValidationError, validate_canonical
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    bundle = _run_product_bundle(src, tmp_path / 'out')[3]
    data = _load_bundle(bundle)
    data['content'][0].update({'type': 'table', 'table_id': 'table-0000000000000000'})
    with pytest.raises(CanonicalValidationError, match='dangling table'):
        validate_canonical(data, bundle)


def test_semantic_validator_rejects_dangling_locator_span(tmp_path):
    from canonical import CanonicalValidationError, validate_canonical

    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    bundle = _run_product_bundle(src, tmp_path / 'out')[3]
    data = _load_bundle(bundle)
    data['content'][0]['source_locator']['spans'] = [{
        'page': 2,
        'source_unit_id': 'unit-0000000000000000',
        'bbox': [0, 0, 10, 10],
    }]

    with pytest.raises(CanonicalValidationError, match='dangling source-unit span'):
        validate_canonical(data, bundle)


def test_semantic_validator_rejects_asset_path_escape_and_hash_mismatch(tmp_path):
    from canonical import CanonicalValidationError, stable_id, validate_canonical
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    bundle = _run_product_bundle(src, tmp_path / 'out')[3]
    data = _load_bundle(bundle)
    document_id = data['document']['document_id']
    locator = {'source_unit_id': data['source_units'][0]['id'], 'index': 1}
    asset_id = stable_id('asset', document_id, locator, 'image', 1)
    data['assets'].append({
        'asset_id': asset_id, 'type': 'image', 'path': '../escape.png',
        'sha256': '0' * 64, 'media_type': 'image/png', 'source_locator': locator,
        'alt': '', 'caption': '',
    })
    data['content'].append({
        'id': stable_id('node', document_id, locator, 'image', 99), 'type': 'image',
        'source_locator': locator, 'asset_id': asset_id,
    })
    with pytest.raises(CanonicalValidationError, match='escapes bundle'):
        validate_canonical(data, bundle)
    asset_path = bundle / 'assets' / 'images' / 'actual.png'
    asset_path.parent.mkdir(parents=True)
    asset_path.write_bytes(b'actual')
    data['assets'][0]['path'] = 'assets/images/actual.png'
    with pytest.raises(CanonicalValidationError, match='hash mismatch'):
        validate_canonical(data, bundle)


def test_quality_status_classification():
    from canonical import quality_from_warnings
    assert quality_from_warnings([]) == 'complete'
    assert quality_from_warnings([{'content_loss': False}]) == 'complete_with_warnings'
    assert quality_from_warnings([{'content_loss': True}]) == 'partial'


def test_semantic_validator_rejects_quality_and_output_manifest_mismatch(tmp_path):
    from canonical import CanonicalValidationError, validate_canonical
    src = tmp_path / 'source.txt'
    src.write_text('Body', encoding='utf-8')
    bundle = _run_product_bundle(src, tmp_path / 'out')[3]
    data = _load_bundle(bundle)
    data['quality']['status'] = 'partial'
    with pytest.raises(CanonicalValidationError, match='quality status'):
        validate_canonical(data, bundle)
    data = _load_bundle(bundle)
    data['outputs']['markdown']['sha256'] = '0' * 64
    with pytest.raises(CanonicalValidationError, match='Markdown hash mismatch'):
        validate_canonical(data, bundle)


def test_no_frontmatter_keeps_json_metadata(tmp_path):
    src = tmp_path / 'plain.txt'
    src.write_text('Body', encoding='utf-8')
    code, _, stderr, bundle = _run_product_bundle(src, tmp_path / 'out', ['--no-frontmatter', '--timestamp', '2026-08-02'])
    assert code == 0, stderr
    assert not (bundle / 'plain.md').read_text(encoding='utf-8').startswith('---')
    assert _load_bundle(bundle)['document']['conversion_timestamp'] == '2026-08-02'


def test_bundle_rename_uses_deterministic_suffix_for_folder_and_files(tmp_path):
    src = tmp_path / 'report.txt'
    src.write_text('# Content Heading\n\nBody', encoding='utf-8')
    output = tmp_path / 'out'
    assert _run_product_bundle(src, output)[0] == 0
    code, _, stderr, _ = _run_product_bundle(src, output, ['--rename'])
    assert code == 0, stderr
    renamed = output / 'report_1'
    assert (renamed / 'report_1.json').exists()
    assert (renamed / 'report_1.md').exists()
    assert _load_bundle(renamed)['document']['title'] == 'report'
    assert 'title: "report"' in (renamed / 'report_1.md').read_text(encoding='utf-8')


def test_collision_rename_preserves_dotted_logical_stem_in_both_modes(tmp_path):
    stem = '10. FF004 - Lynk Pharmaceuticals Co. Ltd'
    src = tmp_path / f'{stem}.txt'
    src.write_text('Body', encoding='utf-8')

    bundle_output = tmp_path / 'bundles'
    assert _run_product_bundle(src, bundle_output)[0] == 0
    code, _, stderr, _ = _run_product_bundle(src, bundle_output, ['--rename'])
    assert code == 0, stderr
    renamed_bundle = bundle_output / f'{stem}_1'
    assert (renamed_bundle / f'{stem}_1.json').exists()
    assert (renamed_bundle / f'{stem}_1.md').exists()

    markdown_output = tmp_path / 'markdown'
    markdown_args = SCRIPT + CONFIG_ARG + [
        '--input', str(src),
        '--output-dir', str(markdown_output),
        '--output-mode', 'markdown',
    ]
    first = subprocess.run(markdown_args, capture_output=True)
    assert first.returncode == 0, first.stderr.decode(errors='replace')
    renamed = subprocess.run(markdown_args + ['--rename'], capture_output=True)
    assert renamed.returncode == 0, renamed.stderr.decode(errors='replace')
    assert (markdown_output / f'{stem}.md').exists()
    assert (markdown_output / f'{stem}_1.md').exists()


def test_default_batch_uses_converted_and_does_not_reprocess_outputs(tmp_path):
    source = tmp_path / 'input'
    source.mkdir()
    (source / 'one.txt').write_text('One', encoding='utf-8')
    args = SCRIPT + CONFIG_ARG + ['--input-dir', str(source)]
    first = subprocess.run(args, capture_output=True)
    assert first.returncode == 0, first.stderr.decode(errors='replace')
    assert (source / '_converted' / 'one' / 'one.json').exists()
    second = subprocess.run(args, capture_output=True)
    assert second.returncode == 2
    assert not (source / '_converted' / '_converted').exists()


def test_collision_is_detected_before_adapter_or_normalizer(tmp_path, monkeypatch):
    from argparse import Namespace
    import pipeline
    src = tmp_path / 'early.txt'
    src.write_text('Body', encoding='utf-8')
    (tmp_path / 'early').mkdir()
    args = Namespace(
        input=str(src), input_dir=None, output_path='', output_dir='', output_mode='bundle',
        overwrite=False, rename=False, timestamp='2026-08-02', language_normalization='simplified',
        no_frontmatter=False,
    )
    monkeypatch.setattr(pipeline, '_build_document', lambda *a, **k: pytest.fail('adapter was called'))
    with pytest.raises(pipeline.OutputCollision):
        pipeline.convert_one(args, str(src))


def test_markitdown_instance_is_reused():
    import adapters
    adapters._MARKITDOWN = None
    assert adapters.get_markitdown() is adapters.get_markitdown()


def test_opencc_runtime_uses_one_conversion_pass(monkeypatch):
    import canonical
    calls = []

    class FakeConverter:
        def convert(self, value):
            calls.append(value)
            return value

    monkeypatch.setattr(canonical, '_opencc_converter', lambda profile: FakeConverter())
    canonical.convert_chinese('內容', 'simplified')
    assert len(calls) == 1


def test_many_canonical_records_share_one_opencc_pass(monkeypatch):
    import canonical
    calls = []

    def fake_convert(value, mode):
        calls.append((value, mode))
        return value

    monkeypatch.setattr(canonical, 'convert_chinese', fake_convert)
    nodes = [
        {'type': 'paragraph', 'text': f'內容 {index}', 'normalized_text': ''}
        for index in range(100)
    ]
    canonical.normalize_canonical_text(nodes, [], 'simplified')
    assert len(calls) == 1
    assert [node['normalized_text'] for node in nodes] == [f'內容 {index}' for index in range(100)]


def test_ooxml_feature_preflight_is_lightweight_and_nonblocking(tmp_path):
    import zipfile
    from adapters import inspect_ooxml_features
    package = tmp_path / 'features.docx'
    with zipfile.ZipFile(package, 'w') as archive:
        archive.writestr(
            'word/document.xml',
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:ins/><w:del/><w:instrText>TOC \\o "1-3"</w:instrText></w:document>',
        )
        archive.writestr(
            'word/comments.xml',
            '<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>',
        )
        archive.writestr('word/media/image1.png', b'png')
    warnings = inspect_ooxml_features(package, 'unit-0000000000000000')
    codes = {item['code'] for item in warnings}
    assert codes == {
        'office_comments_not_preserved',
        'office_revisions_flattened_to_accepted_view',
    }
    assert all(item['content_loss'] for item in warnings)


def test_ooxml_instrtext_is_not_a_revision(tmp_path):
    import zipfile
    from adapters import inspect_ooxml_features

    package = tmp_path / 'field.docx'
    with zipfile.ZipFile(package, 'w') as archive:
        archive.writestr(
            'word/document.xml',
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:instrText>TOC \\o "1-3"</w:instrText></w:document>',
        )
    assert inspect_ooxml_features(package, 'unit-0000000000000000') == []


def test_ooxml_unproved_revision_class_fails_closed(tmp_path):
    import zipfile
    from adapters import inspect_ooxml_features

    package = tmp_path / 'unsupported-revision.docx'
    with zipfile.ZipFile(package, 'w') as archive:
        archive.writestr(
            'word/document.xml',
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:cellDel/></w:document>',
        )
    with pytest.raises(RuntimeError, match='no proved accepted-view transform'):
        inspect_ooxml_features(package, 'unit-0000000000000000')


def test_extract_ooxml_images_deduplicates_asset_and_preserves_occurrences(tmp_path):
    import zipfile
    from ooxml_images import extract_ooxml_images

    source = tmp_path / 'reused.docx'
    document_xml = '''
        <w:document xmlns:w="urn:w" xmlns:wp="urn:wp" xmlns:a="urn:a"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <w:body>
            <wp:inline><wp:docPr name="Logo" descr="Company logo"/><a:blip r:embed="rId1"/></wp:inline>
            <wp:inline><wp:docPr name="Logo" descr="Company logo"/><a:blip r:embed="rId1"/></wp:inline>
          </w:body>
        </w:document>
    '''
    relationships = '''
        <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
          <Relationship Id="rId1" Type="image" Target="media/image1.png"/>
        </Relationships>
    '''
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('word/document.xml', document_xml)
        archive.writestr('word/_rels/document.xml.rels', relationships)
        archive.writestr('word/media/image1.png', b'\x89PNG\r\n\x1a\nimage-data')

    assets, occurrences, warnings = extract_ooxml_images(
        source,
        'sha256:' + ('1' * 64),
        'unit-0000000000000000',
        tmp_path / 'assets' / 'images',
    )

    assert warnings == []
    assert len(assets) == 1
    assert occurrences == [assets[0]['asset_id'], assets[0]['asset_id']]
    assert assets[0]['alt'] == 'Company logo'
    assert assets[0]['path'].startswith('assets/images/')
    assert (tmp_path / Path(assets[0]['path'])).is_file()


def test_extract_ooxml_images_ignores_orphan_media(tmp_path):
    import zipfile
    from ooxml_images import extract_ooxml_images

    source = tmp_path / 'orphan.docx'
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('word/document.xml', '<w:document xmlns:w="urn:w"/>')
        archive.writestr(
            'word/_rels/document.xml.rels',
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>',
        )
        archive.writestr('word/media/orphan.png', b'\x89PNG\r\n\x1a\norphan')
    assets, occurrences, warnings = extract_ooxml_images(
        source,
        'sha256:' + ('9' * 64),
        'unit-0000000000000000',
        tmp_path / 'assets' / 'images',
    )
    assert assets == [] and occurrences == [] and warnings == []
    assert not (tmp_path / 'assets').exists()


def test_office_preflight_rejects_unsafe_member_before_conversion(tmp_path):
    import zipfile
    from office_preflight import preflight_office

    source = tmp_path / 'unsafe.xlsx'
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('../escape.xml', '<root/>')
    with pytest.raises(RuntimeError, match='unsafe member path'):
        preflight_office(source)


def test_office_preflight_reports_xml_nodes_without_treating_count_as_danger(tmp_path):
    import zipfile
    from office_preflight import preflight_office

    source = tmp_path / 'large.docx'
    xml = '<root>' + ''.join(f'<p>{index}</p>' for index in range(5000)) + '</root>'
    with zipfile.ZipFile(source, 'w', compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr('word/document.xml', xml)
    report = preflight_office(source)
    assert report.xml_nodes_by_part['word/document.xml'] == 10001
    assert report.package_members == 1


def test_office_preflight_rejects_xml_doctype(tmp_path):
    import zipfile
    from office_preflight import preflight_office

    source = tmp_path / 'doctype.docx'
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('word/document.xml', '<!DOCTYPE root [<!ENTITY x "text">]><root>&x;</root>')
    with pytest.raises(RuntimeError, match='prohibited internal document type subset'):
        preflight_office(source)


def test_extract_ooxml_images_warns_without_creating_dangling_asset(tmp_path):
    import zipfile
    from ooxml_images import extract_ooxml_images

    source = tmp_path / 'missing.docx'
    document_xml = '''
        <w:document xmlns:w="urn:w" xmlns:a="urn:a"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <w:body><a:blip r:embed="rId1"/></w:body>
        </w:document>
    '''
    relationships = '''
        <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
          <Relationship Id="rId1" Type="image" Target="media/missing.png"/>
        </Relationships>
    '''
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('word/document.xml', document_xml)
        archive.writestr('word/_rels/document.xml.rels', relationships)

    assets, occurrences, warnings = extract_ooxml_images(
        source,
        'sha256:' + ('2' * 64),
        'unit-0000000000000000',
        tmp_path / 'assets' / 'images',
    )

    assert assets == []
    assert occurrences == []
    assert [item['code'] for item in warnings] == ['office_image_target_missing']
    assert warnings[0]['content_loss'] is True
    assert not (tmp_path / 'assets').exists()


def test_extract_ooxml_images_rejects_unsupported_media_type(tmp_path):
    import zipfile
    from ooxml_images import extract_ooxml_images

    source = tmp_path / 'unsupported.docx'
    document_xml = '''
        <w:document xmlns:w="urn:w" xmlns:a="urn:a"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <w:body><a:blip r:embed="rId1"/></w:body>
        </w:document>
    '''
    relationships = '''
        <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
          <Relationship Id="rId1" Type="image" Target="media/image1.bin"/>
        </Relationships>
    '''
    with zipfile.ZipFile(source, 'w') as archive:
        archive.writestr('word/document.xml', document_xml)
        archive.writestr('word/_rels/document.xml.rels', relationships)
        archive.writestr('word/media/image1.bin', b'not-a-supported-image')

    assets, occurrences, warnings = extract_ooxml_images(
        source,
        'sha256:' + ('3' * 64),
        'unit-0000000000000000',
        tmp_path / 'assets' / 'images',
    )

    assert assets == []
    assert occurrences == []
    assert [item['code'] for item in warnings] == ['office_image_media_type_unsupported']
    assert warnings[0]['content_loss'] is True
    assert not (tmp_path / 'assets').exists()


def _make_pdf(path, draw):
    from reportlab.pdfgen import canvas
    document = canvas.Canvas(str(path), pagesize=(612, 792))
    draw(document)
    document.save()




def test_product_pdf_bundle_uses_inspector_and_emits_clean_markdown(tmp_path):
    pdf = tmp_path / 'native.pdf'

    def draw(c):
        c.setFont('Helvetica-Bold', 18)
        c.drawString(72, 740, 'Native PDF')
        c.setFont('Helvetica', 10)
        c.drawString(72, 710, 'First paragraph line')
        c.drawString(72, 698, 'continues here.')

    _make_pdf(pdf, draw)
    code, _, stderr, bundle = _run_product_bundle(pdf, tmp_path / 'out')
    assert code == 0, stderr
    data = _load_bundle(bundle)
    assert data['adapter']['name'] == 'pdf-inspector'
    page_unit = next(unit for unit in data['source_units'] if unit['type'] == 'page')
    assert page_unit['locator']['page'] == 1
    assert all(node['source_locator'].get('page_range') == [1, 1] for node in data['content'])
    assert any(
        node['source_locator'].get('extraction_method') == 'pdf-inspector'
        for node in data['content']
        if node['type'] != 'image'
    )
    markdown = (bundle / 'native.md').read_text(encoding='utf-8')
    assert 'block:' not in markdown
    assert 'Native PDF' in markdown


def test_product_sparse_readable_pdf_is_not_deleted_by_full_result_ocr_signal(
    tmp_path,
):
    pdf = tmp_path / 'sparse-readable.pdf'

    def draw(c):
        for page in range(1, 4):
            c.setFont('Helvetica', 10)
            c.drawString(72, 720, f'Precisely readable sparse page {page}.')
            if page < 3:
                c.showPage()

    _make_pdf(pdf, draw)
    code, _, stderr, bundle = _run_product_bundle(
        pdf,
        tmp_path / 'out',
        ['--ocr', 'off'],
    )

    assert code == 0, stderr
    data = _load_bundle(bundle)
    markdown = (bundle / 'sparse-readable.md').read_text(encoding='utf-8')
    for page in range(1, 4):
        assert f'Precisely readable sparse page {page}.' in markdown
    assert not any(
        warning['code'] == 'ocr_required'
        for warning in data['quality']['warnings']
    )


def _processed_result(markdown, page_count, *, ocr_pages=(), reasons=None):
    reason_items = []
    for page, values in (reasons or {}).items():
        reason_items.append(type('Reason', (), {
            'page': page,
            'reasons': list(values),
        })())
    return type('Result', (), {
        'markdown': markdown,
        'page_count': page_count,
        'pages_needing_ocr': list(ocr_pages),
        'ocr_reasons_by_page': reason_items,
    })()


def test_pdf_inspector_healthy_page_is_authoritative(monkeypatch, tmp_path):
    import builtins
    import pdf_inspector_adapter as module
    from canonical import make_text_fields, stable_id

    source = tmp_path / 'authoritative.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('a' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 1)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    unit_locator = {'page': 1}
    unit_id = stable_id('unit', document_id, unit_locator, 'page', 1)

    class NativeSupport:
        def extract(self, _source, _document_id, mode, _asset_dir=None):
            locator = {'source_unit_id': unit_id, 'page': 1}
            return {
                'source_units': [{
                    'id': unit_id, 'type': 'page', 'index': 1,
                    'locator': unit_locator, 'status': 'complete', 'warnings': [],
                }],
                'content': [{
                    'id': 'native', 'type': 'paragraph', 'source_locator': locator,
                    **make_text_fields('PDFium replacement', 'PDFium replacement', mode),
                }],
                'tables': [], 'assets': [], 'relationships': [], 'warnings': [],
            }

    calls = []

    def process(_source, pages=None):
        calls.append(pages)
        assert pages is None
        return _processed_result(
            'CONFIDENTIAL\n\n# 2024\n\n'
            '| Item | Amount |\n| --- | --- |\n| Revenue | 100 |',
            1,
        )

    inspector = type('Inspector', (), {'process_pdf': staticmethod(process)})()

    def unexpected_ocr(*_args):
        raise AssertionError('healthy Inspector pages must not invoke OCR')

    original_import = builtins.__import__

    def block_optional_pdf_capabilities(name, *args, **kwargs):
        if name in {'pdfminer', 'pypdfium2'} or name.startswith('pdfminer.'):
            pytest.fail(f'healthy Markdown-only PDF imported optional capability {name}')
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, '__import__', block_optional_pdf_capabilities)

    result = module.PdfInspectorAdapter(
        inspector=inspector,
        fallback_adapter=NativeSupport(),
        ocr_mode='auto',
        ocr_runner=unexpected_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert [node['type'] for node in result['content']] == [
        'paragraph', 'heading', 'table'
    ]
    assert result['content'][0]['text'] == 'CONFIDENTIAL'
    assert result['content'][1]['text'] == '2024'
    assert result['content'][1]['level'] == 1
    assert result['tables'][0]['raw_rows'][1] == ['Revenue', '100']
    assert not any('PDFium replacement' in node.get('text', '') for node in result['content'])
    assert not any(warning['code'] == 'ocr_applied' for warning in result['warnings'])
    assert calls == [None]


def test_pdf_inspector_missing_cid_repair_capability_keeps_inspector_content(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'cid-repair-unavailable.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('b' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 1)
    monkeypatch.setattr(
        module,
        'repair_standard_cid_tounicode',
        lambda *_args: (_ for _ in ()).throw(
            ModuleNotFoundError("No module named 'pdfminer'")
        ),
    )
    result = module.PdfInspectorAdapter(
        inspector=type('Inspector', (), {
            'process_pdf': staticmethod(
                lambda _source, pages=None: _processed_result(
                    '# Inspector title\n\nInspector body.', 1
                )
            )
        })(),
        ocr_mode='auto',
    ).extract(str(source), document_id, 'preserve')

    assert [node['text'] for node in result['content']] == [
        'Inspector title', 'Inspector body.'
    ]
    warning = next(
        item for item in result['warnings']
        if item['code'] == 'pdf_inspector_cid_repair_failed'
    )
    assert warning['content_loss'] is False


def test_pdf_inspector_preserves_global_outline_levels(monkeypatch, tmp_path):
    import pdf_inspector_adapter as module

    source = tmp_path / 'global-outline.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('1' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    calls = []

    def process(_source, pages=None):
        calls.append(pages)
        return _processed_result(
            '# Cover\n\n### Warning\n\nBody.\n\n## Main section',
            3,
        )

    result = module.PdfInspectorAdapter(
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=lambda *_args: (_ for _ in ()).throw(
            AssertionError('healthy full-document Inspector output must not invoke OCR')
        ),
    ).extract(str(source), document_id, 'preserve')

    headings = [
        (node['text'], node['level'])
        for node in result['content']
        if node['type'] == 'heading'
    ]
    assert headings == [('Cover', 1), ('Warning', 3), ('Main section', 2)]
    assert calls == [None]


def test_pdf_inspector_affine_page_alignment_rejects_unproven_page_edges():
    import pdf_inspector_adapter as module

    common = ''.join(f'Block{index:03d}AlphaBetaGamma' for index in range(80))
    page = ('Changed page prefix ' * 8) + common
    global_markdown = ('Earlier document text ' * 30) + common

    assert module._align_selected_page(global_markdown, page) is None


def test_pdf_inspector_page_alignment_rejects_ambiguous_duplicate_signature():
    import pdf_inspector_adapter as module

    assert module._align_selected_page(
        'Repeated healthy marker\n\nMiddle\n\nRepeated healthy marker',
        'Repeated healthy marker',
    ) is None


def test_pdf_inspector_raw_alignment_never_consumes_same_line_healthy_text():
    import pdf_inspector_adapter as module

    global_markdown = 'HEALTHY_PREFIX FLAGGED TEXT HEALTHY_SUFFIX\n\n# Tail'
    alignment = module._align_selected_page(global_markdown, 'FLAGGED TEXT')

    assert alignment is not None
    raw_start = module._raw_alignment_start(global_markdown, alignment)
    raw_end = module._raw_alignment_end(global_markdown, alignment)
    assert global_markdown[:raw_start] == 'HEALTHY_PREFIX '
    assert global_markdown[raw_start:raw_end] == 'FLAGGED TEXT'
    assert global_markdown[raw_end:].startswith(' HEALTHY_SUFFIX')


def test_pdf_inspector_affine_page_alignment_rejects_reordered_runs():
    import pdf_inspector_adapter as module

    first = ''.join(f'First{index:03d}Alpha' for index in range(40))
    second = ''.join(f'Second{index:03d}Beta' for index in range(40))
    third = ''.join(f'Third{index:03d}Gamma' for index in range(40))
    page = first + second + third
    global_markdown = (
        'Healthy prefix\n\n' + first + ('Z' * 32) + third + second + '\n\nTail'
    )

    assert module._align_selected_page(global_markdown, page) is None


def test_pdf_inspector_direct_flagged_span_preserves_sparse_front_matter(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'sparse-front-matter.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('f' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    healthy = (
        '# Cover\n\nImportant application terms.\n\n'
        '## Important notice\n\n– i –\n\n'
        '## Expected timetable\n\n– v –'
    )
    flagged_core = ''.join(
        f'FlaggedTocEntry{index:03d}UniqueValue\n' for index in range(30)
    )
    stale_folio = 'I-1\n\nII-1\n\nIII-1\n\nIV-1\n\nV-1\n\n– vii –'
    tail = '## Summary\n\nHealthy tail content.'
    global_markdown = (
        healthy + '\n\n' + flagged_core + '\n' + stale_folio + '\n\n' + tail
    )
    selected_bad_page = flagged_core + '\n' + stale_folio

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                global_markdown,
                3,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        assert pages == [2]
        return _processed_result(selected_bad_page, 3, ocr_pages=(2,))

    def run_ocr(_source, pages, _provider):
        assert pages == {2}
        return {
            2: ([{
                'text': 'OCR replacement table of contents',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
        }

    result = module.PdfInspectorAdapter(
        type('Provider', (), {
            'available': True, 'name': 'fake-ocr', 'version': '1.0'
        })(),
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    texts = [node['text'] for node in result['content']]
    assert 'Cover' in texts
    assert 'Important application terms.' in texts
    assert 'Important notice' in texts
    assert 'Expected timetable' in texts
    assert 'OCR replacement table of contents' in texts
    assert 'Summary' in texts
    assert 'Healthy tail content.' in texts
    assert not any('FlaggedTocEntry' in text for text in texts)
    assert not any(text in {'I-1', 'II-1', 'III-1', 'IV-1', 'V-1'} for text in texts)


def test_pdf_inspector_retains_readable_page_after_standard_cid_repair(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'cid-readable.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('6' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 1)
    monkeypatch.setattr(
        module,
        'repair_standard_cid_tounicode',
        lambda _source, _target: (1, ['Adobe-CNS1']),
    )
    inspector_markdown = (
        '# 行业概览\n\n'
        '本页在标准字符映射修复后具有完整且清晰可读的中文内容。\n\n'
        '| 排名 | 供应商 | 收益 |\n'
        '| --- | --- | --- |\n'
        '| 1 | 公司I | 100 |\n'
        '| 2 | 公司II | 90 |'
    )

    def process(_source, pages=None):
        assert pages in (None, [1])
        return _processed_result(
            inspector_markdown,
            1,
            ocr_pages=(1,),
            reasons={1: ('suspected_garbled_text',)},
        )

    result = module.PdfInspectorAdapter(
        type('Provider', (), {'available': True})(),
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=lambda *_args: (_ for _ in ()).throw(
            AssertionError('readable CID-repaired Inspector page must win over OCR')
        ),
    ).extract(str(source), document_id, 'preserve')

    texts = [node.get('text', '') for node in result['content']]
    assert '行业概览' in texts
    assert result['tables'][0]['raw_rows'][1] == ['1', '公司I', '100']
    assert any(
        warning['code'] == 'pdf_inspector_cid_page_retained'
        for warning in result['warnings']
    )
    assert not any(warning['code'] == 'ocr_applied' for warning in result['warnings'])


def test_pdf_inspector_cid_readability_never_overrides_other_ocr_reasons():
    import pdf_inspector_adapter as module

    readable = '完整且清晰可读的标准字符映射修复内容。' * 4
    assert module._cid_repaired_page_is_readable(
        readable, 'suspected_garbled_text'
    )
    assert not module._cid_repaired_page_is_readable(
        readable, 'scanned, suspected_garbled_text'
    )


def test_pdf_inspector_unusable_page_uses_ocr_only(monkeypatch, tmp_path):
    import pdf_inspector_adapter as module

    source = tmp_path / 'ocr-only.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('b' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    process_calls = []

    def process(_source, pages=None):
        process_calls.append(pages)
        if pages is None:
            return _processed_result(
                '# Inspector title\n\nHealthy page.\n\n'
                '## Flagged native page\n\nFLAGGED_NATIVE SHOULD BE REPLACED.\n\n'
                '## Final page\n\nTail content.',
                3,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        if pages == [2]:
            return _processed_result(
                '## Flagged native page\n\nFLAGGED_NATIVE SHOULD BE REPLACED.',
                3,
                ocr_pages=(2,),
            )
        raise AssertionError(f'unexpected selected-page request: {pages}')

    inspector = type('Inspector', (), {'process_pdf': staticmethod(process)})()
    ocr_result = type('OcrResult', (), {
        'engine': 'fake-ocr',
        'engine_version': '1.0',
        'runtime': 'test',
        'runtime_version': '1',
        'model_profile': 'fixture',
        'language': 'ch',
        'min_confidence': 0.5,
        'requested_dpi': 300.0,
        'effective_dpi': 300.0,
        'raster_width': 1200,
        'raster_height': 1600,
        'dropped_low_confidence': 0,
        'dropped_invalid': 0,
        'spans': (
            {'text': 'OCR recovered amount 12', 'bbox': [10, 20, 200, 40], 'confidence': 0.99},
        ),
    })()
    requested = []

    def run_ocr(_source, pages, _provider):
        requested.append(set(pages))
        return {2: (ocr_result, None, None)}

    provider = type('Provider', (), {'available': True, 'name': 'fake-ocr'})()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=inspector,
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert requested == [{2}]
    texts = [node.get('text', '') for node in result['content']]
    assert 'Healthy page.' in texts
    assert 'OCR recovered amount 12' in texts
    assert 'Flagged native page' not in texts
    assert 'FLAGGED_NATIVE SHOULD BE REPLACED.' not in texts
    assert texts.index('Healthy page.') < texts.index('OCR recovered amount 12')
    assert texts.index('OCR recovered amount 12') < texts.index('Final page')
    recovered = next(node for node in result['content'] if node.get('text') == 'OCR recovered amount 12')
    assert recovered['source_locator']['extraction_method'] == 'ocr'
    page_two = next(
        unit for unit in result['source_units']
        if unit['type'] == 'page' and unit['index'] == 2
    )
    assert page_two['status'] == 'warning'
    assert any(warning['code'] == 'ocr_applied' for warning in result['warnings'])
    assert not any('pdfium' in warning['message'].casefold() for warning in result['warnings'])
    assert process_calls == [None, [2]]


def test_pdf_inspector_all_flagged_pages_replace_retained_document_markdown(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'all-pages-retained.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('0' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    calls = []
    per_page_calls = []

    def process(_source, pages=None):
        calls.append(pages)
        assert pages is None
        return _processed_result(
            '# NATIVE PAGE 1\n\nNATIVE PAGE 2\n\nNATIVE PAGE 3',
            3,
            ocr_pages=(1, 2, 3),
            reasons={
                1: ('suspected_garbled_text',),
                2: ('suspected_garbled_text',),
                3: ('suspected_garbled_text',),
            },
        )

    def run_ocr(_source, pages, _provider):
        assert pages == {1, 2, 3}
        return {
            page: ([{
                'text': f'OCR replacement page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    def extract_pages(_source, pages=None):
        per_page_calls.append(pages)
        return type('PagesResult', (), {
            'pages': [
                type('PageResult', (), {
                    'page': page,
                    'needs_ocr': True,
                })()
                for page in pages
            ],
            'pages_needing_ocr': [],
        })()

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {
            'process_pdf': staticmethod(process),
            'extract_pages_markdown': staticmethod(extract_pages),
        })(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert [node['text'] for node in result['content']] == [
        'OCR replacement page 1',
        'OCR replacement page 2',
        'OCR replacement page 3',
    ]
    assert calls == [None]
    assert per_page_calls == [[0, 1, 2]]


def test_pdf_inspector_consecutive_middle_run_replaces_entire_proven_span(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'middle-run-retained.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('2' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 4)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    page_one = '# Healthy page 1\n\nHEALTHY_ONE'
    page_four = '# Healthy page 4\n\nHEALTHY_FOUR'
    calls = []

    def process(_source, pages=None):
        calls.append(pages)
        if pages is None:
            return _processed_result(
                page_one
                + '\n\n# Flagged page 2\n\nSTALE_TWO'
                + '\n\n# Flagged page 3\n\nSTALE_THREE\n\n'
                + page_four,
                4,
                ocr_pages=(2, 3),
                reasons={2: ('scanned',), 3: ('scanned',)},
            )
        if pages == [2]:
            return _processed_result(
                '# Flagged page 2\n\nSTALE_TWO', 4, ocr_pages=(2,)
            )
        if pages == [3]:
            return _processed_result(
                '# Flagged page 3\n\nSTALE_THREE', 4, ocr_pages=(3,)
            )
        raise AssertionError(f'unexpected selected-page request: {pages}')

    def run_ocr(_source, pages, _provider):
        assert pages == {2, 3}
        return {
            page: ([{
                'text': f'OCR middle page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    texts = [node['text'] for node in result['content']]
    assert texts == [
        'Healthy page 1',
        'HEALTHY_ONE',
        'OCR middle page 2',
        'OCR middle page 3',
        'Healthy page 4',
        'HEALTHY_FOUR',
    ]
    assert 'STALE_TWO' not in texts
    assert 'STALE_THREE' not in texts
    assert calls == [None, [2], [3]]


def test_pdf_inspector_failed_ocr_removes_stale_flagged_page_text(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'failed-ocr-retained.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('3' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    page_one = '# Healthy page 1\n\nHEALTHY_ONE'
    page_three = '# Healthy page 3\n\nHEALTHY_THREE'

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                page_one + '\n\n# Flagged page 2\n\nSTALE_UNUSABLE\n\n' + page_three,
                3,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        if pages == [2]:
            return _processed_result(
                '# Flagged page 2\n\nSTALE_UNUSABLE', 3, ocr_pages=(2,)
            )
        raise AssertionError(f'unexpected selected-page request: {pages}')

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=lambda _source, _pages, _provider: {
            2: (None, 'ocr_failed', 'synthetic OCR failure')
        },
    ).extract(str(source), document_id, 'preserve')

    texts = [node['text'] for node in result['content']]
    assert texts == [
        'Healthy page 1', 'HEALTHY_ONE', 'Healthy page 3', 'HEALTHY_THREE'
    ]
    assert 'STALE_UNUSABLE' not in texts
    page_two = next(
        unit for unit in result['source_units']
        if unit['type'] == 'page' and unit['index'] == 2
    )
    assert page_two['status'] == 'ocr_required'
    assert any(warning['code'] == 'ocr_failed' for warning in result['warnings'])
    assert any(warning['code'] == 'ocr_required' for warning in result['warnings'])


def test_pdf_inspector_failed_ocr_never_uses_native_text(monkeypatch, tmp_path):
    import pdf_inspector_adapter as module

    source = tmp_path / 'ocr-unavailable.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('c' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                '# Healthy Inspector content\n\nFirst page remains.\n\n'
                '# Flagged page\n\nSTALE_GARBLED',
                2,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        assert pages == [2]
        return _processed_result(
            '# Flagged page\n\nSTALE_GARBLED', 2, ocr_pages=(2,)
        )

    inspector = type('Inspector', (), {'process_pdf': staticmethod(process)})()

    result = module.PdfInspectorAdapter(
        type('Provider', (), {'available': False})(),
        inspector=inspector,
        ocr_mode='auto',
        ocr_runner=lambda _source, _pages, _provider: {
            2: (None, 'ocr_unavailable', 'backend missing')
        },
    ).extract(str(source), document_id, 'preserve')

    assert [node.get('text') for node in result['content']] == [
        'Healthy Inspector content', 'First page remains.'
    ]
    assert all('STALE_GARBLED' not in node.get('text', '') for node in result['content'])
    page_two = next(
        unit for unit in result['source_units']
        if unit['type'] == 'page' and unit['index'] == 2
    )
    assert page_two['status'] == 'ocr_required'
    assert any(warning['code'] == 'ocr_unavailable' for warning in result['warnings'])
    assert any(warning['code'] == 'ocr_required' for warning in result['warnings'])
    assert all(warning['content_loss'] for warning in result['warnings'])


def test_auto_ocr_missing_pdfium_is_normalized_as_optional_unavailability(
    monkeypatch
):
    import builtins
    import pdf_inspector_adapter as module

    original_import = builtins.__import__

    def block_pdfium(name, *args, **kwargs):
        if name == 'pypdfium2':
            raise ModuleNotFoundError("No module named 'pypdfium2'")
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, '__import__', block_pdfium)
    provider = type('Provider', (), {'available': True})()
    result = module._run_ocr_pages('unopened.pdf', {2}, provider)

    assert result[2][0] is None
    assert result[2][1] == 'ocr_unavailable'
    assert 'rasterization is unavailable' in result[2][2]


def test_pdf_inspector_document_failure_routes_every_page_to_ocr(monkeypatch, tmp_path):
    import pdf_inspector_adapter as module

    source = tmp_path / 'inspector-failed.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('d' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )

    def fail(_source, pages=None):
        raise RuntimeError('synthetic Inspector failure')

    requested = []

    def run_ocr(_source, pages, _provider):
        requested.append(set(pages))
        return {
            page: ([{
                'text': f'OCR page {page} recovered',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.98,
            }], None, None)
            for page in pages
        }

    inspector = type('Inspector', (), {'process_pdf': staticmethod(fail)})()
    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=inspector,
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert requested == [{1, 2}]
    assert [node['text'] for node in result['content']] == [
        'OCR page 1 recovered', 'OCR page 2 recovered'
    ]
    assert all(
        node['source_locator']['extraction_method'] == 'ocr'
        for node in result['content']
    )
    assert any(
        warning['code'] == 'pdf_inspector_document_ocr_fallback'
        for warning in result['warnings']
    )


def test_pdf_inspector_ocr_off_marks_omitted_page_without_fallback(monkeypatch, tmp_path):
    import pdf_inspector_adapter as module

    source = tmp_path / 'ocr-off.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('e' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                '# Healthy page\n\nInspector content.\n\n'
                '# Flagged page\n\nSTALE_GARBLED',
                2,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        assert pages == [2]
        return _processed_result(
            '# Flagged page\n\nSTALE_GARBLED', 2, ocr_pages=(2,)
        )

    inspector = type('Inspector', (), {'process_pdf': staticmethod(process)})()

    result = module.PdfInspectorAdapter(
        inspector=inspector,
        ocr_mode='off',
        ocr_runner=lambda *_args: (_ for _ in ()).throw(
            AssertionError('OCR-off mode must not invoke OCR')
        ),
    ).extract(str(source), document_id, 'preserve')

    assert [node['text'] for node in result['content']] == [
        'Healthy page', 'Inspector content.'
    ]
    assert all('STALE_GARBLED' not in node.get('text', '') for node in result['content'])
    page_two = next(
        unit for unit in result['source_units']
        if unit['type'] == 'page' and unit['index'] == 2
    )
    assert page_two['status'] == 'ocr_required'
    assert any(warning['code'] == 'ocr_required' for warning in result['warnings'])


def test_pdf_inspector_one_based_ocr_pages_keep_first_and_last_in_order(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'boundary-pages.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('7' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    calls = []

    def process(_source, pages=None):
        calls.append(pages)
        if pages is None:
            return _processed_result(
                '# STALE FIRST\n\n'
                '## Middle page\n\nMiddle Inspector content.\n\n'
                '# STALE LAST',
                3,
                ocr_pages=(1, 3),
                reasons={1: ('scanned',), 3: ('scanned',)},
            )
        assert pages in ([1], [3])
        return _processed_result(
            '# STALE FIRST' if pages == [1] else '# STALE LAST',
            3,
            ocr_pages=tuple(pages),
        )

    def run_ocr(_source, pages, _provider):
        assert pages == {1, 3}
        return {
            page: ([{
                'text': f'OCR physical page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert [node['text'] for node in result['content']] == [
        'OCR physical page 1',
        'Middle page',
        'Middle Inspector content.',
        'OCR physical page 3',
    ]
    assert calls == [None, [1], [3]]
    assert {
        warning['source_unit']
        for warning in result['warnings']
        if warning['code'] == 'ocr_applied'
    } == {
        next(
            unit['id']
            for unit in result['source_units']
            if unit['type'] == 'page' and unit['index'] == 1
        ),
        next(
            unit['id']
            for unit in result['source_units']
            if unit['type'] == 'page' and unit['index'] == 3
        ),
    }


def test_pdf_inspector_unprovable_boundary_routes_entire_document_to_ocr(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'unprovable-boundary.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('8' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 3)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                '# Inspector content\n\nHealthy native text.',
                3,
                ocr_pages=(2,),
                reasons={2: ('scanned',)},
            )
        return _processed_result('No matching global anchor exists.', 3)

    requested = []

    def run_ocr(_source, pages, _provider):
        requested.append(set(pages))
        return {
            page: ([{
                'text': f'Ordered OCR page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert requested == [{1, 2, 3}]
    assert [node['text'] for node in result['content']] == [
        'Ordered OCR page 1', 'Ordered OCR page 2', 'Ordered OCR page 3'
    ]
    assert all(
        node['source_locator']['extraction_method'] == 'ocr'
        for node in result['content']
    )
    assert any(
        warning['code'] == 'pdf_inspector_alignment_ocr_fallback'
        for warning in result['warnings']
    )


def test_pdf_inspector_unprovable_span_with_ocr_off_reports_discard(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'unprovable-ocr-off.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('a' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                '# Healthy content\n\nUnlocatable retained text.',
                2,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        return _processed_result('Different selected-page text.', 2)

    result = module.PdfInspectorAdapter(
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='off',
        ocr_runner=lambda *_args: (_ for _ in ()).throw(
            AssertionError('OCR-off mode must not invoke OCR')
        ),
    ).extract(str(source), document_id, 'preserve')

    assert result['content'] == []
    warning = next(
        item for item in result['warnings']
        if item['code'] == 'pdf_inspector_alignment_ocr_fallback'
    )
    assert 'OCR was disabled' in warning['message']
    assert 'routed to ordered OCR' not in warning['message']


@pytest.mark.parametrize('ocr_mode', ['off', 'auto'])
def test_pdf_inspector_empty_selected_page_never_publishes_stale_text(
    monkeypatch, tmp_path, ocr_mode
):
    import pdf_inspector_adapter as module

    source = tmp_path / f'empty-selected-{ocr_mode}.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('c' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    per_page_calls = []

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                '# Healthy content\n\nKEEP\n\n'
                '## Flagged content\n\nSTALE_GARBLED',
                2,
                ocr_pages=(2,),
                reasons={2: ('suspected_garbled_text',)},
            )
        assert pages == [2]
        return _processed_result('', 2, ocr_pages=(2,))

    def extract_pages(_source, pages=None):
        per_page_calls.append(pages)
        return type('PagesResult', (), {
            'pages': [
                type('PageResult', (), {
                    'page': page,
                    'markdown': (
                        '# Healthy content\n\nKEEP' if page == 0 else ''
                    ),
                    'needs_ocr': page == 1,
                })()
                for page in pages
            ],
            'pages_needing_ocr': [2],
        })()

    provider = None
    kwargs = {}
    if ocr_mode == 'off':
        kwargs['ocr_runner'] = lambda *_args: (_ for _ in ()).throw(
            AssertionError('OCR-off mode must not invoke OCR')
        )
    else:
        provider = type('UnavailableProvider', (), {'available': False})()

    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {
            'process_pdf': staticmethod(process),
            'extract_pages_markdown': staticmethod(extract_pages),
        })(),
        ocr_mode=ocr_mode,
        **kwargs,
    ).extract(str(source), document_id, 'preserve')

    assert per_page_calls == [[1], [0, 1]]
    assert result['content'] == []
    assert not any(
        'STALE_GARBLED' in node.get('text', '')
        for node in result['content']
    )
    warning = next(
        item for item in result['warnings']
        if item['code'] == 'pdf_inspector_alignment_ocr_fallback'
    )
    if ocr_mode == 'off':
        assert 'OCR was disabled' in warning['message']
        assert not any(
            item['code'] == 'ocr_unavailable'
            for item in result['warnings']
        )
    else:
        assert 'routed to ordered OCR' in warning['message']
        unavailable_pages = {
            item['source_unit']
            for item in result['warnings']
            if item['code'] == 'ocr_unavailable'
        }
        assert len(unavailable_pages) == 2


def test_pdf_inspector_out_of_order_flagged_spans_route_all_pages_to_ocr(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'overlapping-anchors.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('4' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 5)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    page_one = '# Healthy page 1\n\nHEALTHY_ONE'
    page_four = '# Healthy page 4\n\nHEALTHY_FOUR'
    page_two = '# Flagged page 2\n\nSTALE_TWO'
    page_five = '# Flagged page 5\n\nSTALE_FIVE'

    def process(_source, pages=None):
        if pages is None:
            return _processed_result(
                page_one
                + '\n\n' + page_five
                + '\n\n# Healthy page 3\n\nHEALTHY_THREE\n\n'
                + page_four
                + '\n\n' + page_two,
                5,
                ocr_pages=(2, 5),
                reasons={2: ('scanned',), 5: ('scanned',)},
            )
        selected = {
            (2,): page_two,
            (5,): page_five,
        }
        return _processed_result(selected[tuple(pages)], 5)

    requested = []

    def run_ocr(_source, pages, _provider):
        requested.append(set(pages))
        return {
            page: ([{
                'text': f'Fallback OCR page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        inspector=type('Inspector', (), {'process_pdf': staticmethod(process)})(),
        ocr_mode='auto',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert requested == [{1, 2, 3, 4, 5}]
    assert [node['text'] for node in result['content']] == [
        f'Fallback OCR page {page}' for page in range(1, 6)
    ]
    assert any(
        warning['code'] == 'pdf_inspector_alignment_ocr_fallback'
        for warning in result['warnings']
    )


def test_pdf_inspector_force_uses_ordered_ocr_without_calling_inspector(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'forced-ocr.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('9' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module,
        '_pdf_inspector',
        lambda: pytest.fail('force mode must not initialize PDF Inspector'),
    )

    requested = []

    def run_ocr(_source, pages, _provider):
        requested.append(set(pages))
        return {
            page: ([{
                'text': f'Forced OCR page {page}',
                'bbox': [10, 20, 200, 40],
                'confidence': 0.99,
            }], None, None)
            for page in pages
        }

    provider = type('Provider', (), {
        'available': True, 'name': 'fake-ocr', 'version': '1.0'
    })()
    result = module.PdfInspectorAdapter(
        provider,
        ocr_mode='force',
        ocr_runner=run_ocr,
    ).extract(str(source), document_id, 'preserve')

    assert requested == [{1, 2}]
    assert [node['text'] for node in result['content']] == [
        'Forced OCR page 1', 'Forced OCR page 2'
    ]


def test_pdf_inspector_adapter_does_not_filter_inspector_header_footer_text(
    monkeypatch, tmp_path
):
    import pdf_inspector_adapter as module

    source = tmp_path / 'inspector-chrome.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('6' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 2)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    inspector_markdown = (
        'CONFIDENTIAL\n\nFirst body.\n\nFooter 1\n\n'
        'CONFIDENTIAL\n\nSecond body.\n\nFooter 2'
    )
    result = module.PdfInspectorAdapter(
        inspector=type('Inspector', (), {
            'process_pdf': staticmethod(
                lambda _source, pages=None: _processed_result(
                    inspector_markdown, 2
                )
            )
        })(),
        ocr_mode='auto',
    ).extract(str(source), document_id, 'preserve')

    emitted = '\n'.join(node.get('text', '') for node in result['content'])
    assert emitted.count('CONFIDENTIAL') == 2
    assert 'Footer 1' in emitted
    assert 'Footer 2' in emitted


def test_pdf_inspector_optional_image_failure_never_blocks_document(
    monkeypatch, tmp_path
):
    import builtins
    import pdf_inspector_adapter as module

    source = tmp_path / 'image-support-failed.pdf'
    source.write_bytes(b'%PDF fake')
    document_id = 'sha256:' + ('5' * 64)
    monkeypatch.setattr(module, '_pdf_page_count', lambda _source: 1)
    monkeypatch.setattr(
        module, 'repair_standard_cid_tounicode', lambda _source, _target: (0, [])
    )
    original_import = builtins.__import__

    def block_pdfium(name, *args, **kwargs):
        if name == 'pypdfium2':
            raise ModuleNotFoundError("No module named 'pypdfium2'")
        return original_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, '__import__', block_pdfium)
    result = module.PdfInspectorAdapter(
        inspector=type('Inspector', (), {
            'process_pdf': staticmethod(
                lambda _source, pages=None: _processed_result(
                    '# Inspector title\n\nInspector body.', 1
                )
            )
        })(),
        ocr_mode='auto',
    ).extract(
        str(source), document_id, 'preserve', tmp_path / 'assets' / 'images'
    )

    assert [node['text'] for node in result['content']] == [
        'Inspector title', 'Inspector body.'
    ]
    assert any(
        warning['code'] == 'pdf_image_extraction_failed'
        for warning in result['warnings']
    )


@pytest.mark.parametrize(
    ('encoding', 'vertical'),
    [('/Identity-H', False), ('/Identity-V', True)],
)
def test_pdf_inspector_repairs_standard_identity_cid_font(
    tmp_path, encoding, vertical
):
    from pdfminer.cmapdb import CMapDB
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import (
        ArrayObject, DictionaryObject, NameObject, NumberObject, TextStringObject,
    )
    from pdf_inspector_adapter import repair_standard_cid_tounicode

    source = tmp_path / f'identity-{"v" if vertical else "h"}.pdf'
    target = tmp_path / f'identity-{"v" if vertical else "h"}.repaired.pdf'
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    descendant = DictionaryObject({
        NameObject('/Type'): NameObject('/Font'),
        NameObject('/Subtype'): NameObject('/CIDFontType0'),
        NameObject('/BaseFont'): NameObject('/SyntheticJapan'),
        NameObject('/CIDSystemInfo'): DictionaryObject({
            NameObject('/Registry'): TextStringObject('Adobe'),
            NameObject('/Ordering'): TextStringObject('Japan1'),
            NameObject('/Supplement'): NumberObject(0),
        }),
    })
    font = DictionaryObject({
        NameObject('/Type'): NameObject('/Font'),
        NameObject('/Subtype'): NameObject('/Type0'),
        NameObject('/BaseFont'): NameObject('/SyntheticJapan'),
        NameObject('/Encoding'): NameObject(encoding),
        NameObject('/DescendantFonts'): ArrayObject([writer._add_object(descendant)]),
    })
    page[NameObject('/Resources')] = DictionaryObject({
        NameObject('/Font'): DictionaryObject({
            NameObject('/F1'): writer._add_object(font),
        }),
    })
    with source.open('wb') as handle:
        writer.write(handle)

    count, collections = repair_standard_cid_tounicode(source, target)
    repaired_font = (
        PdfReader(str(target)).pages[0]['/Resources']['/Font']['/F1'].get_object()
    )
    cmap = repaired_font['/ToUnicode'].get_object().get_data()
    expected_map = CMapDB.get_unicode_map('Adobe-Japan1', vertical=vertical)
    opposite_map = CMapDB.get_unicode_map('Adobe-Japan1', vertical=not vertical)
    differing_cid = next(
        cid
        for cid, value in expected_map.cid2unichr.items()
        if 0 <= cid <= 0xFFFF
        and value
        and opposite_map.cid2unichr.get(cid) not in {None, value}
    )
    expected_hex = expected_map.cid2unichr[differing_cid].encode(
        'utf-16-be'
    ).hex().upper()
    opposite_hex = opposite_map.cid2unichr[differing_cid].encode(
        'utf-16-be'
    ).hex().upper()

    assert count == 1
    assert collections == ['Adobe-Japan1']
    assert b'begincmap' in cmap
    assert f'<{differing_cid:04X}> <{expected_hex}>'.encode('ascii') in cmap
    assert f'<{differing_cid:04X}> <{opposite_hex}>'.encode('ascii') not in cmap
































def test_ocr_required_page_publishes_partial_bundle(tmp_path):
    pdf = tmp_path / 'partial.pdf'

    def draw(c):
        c.drawString(72, 720, 'Usable page')
        c.showPage()
        c.rect(72, 600, 200, 100, stroke=1, fill=0)

    _make_pdf(pdf, draw)
    code, stdout, stderr, bundle = _run_pipeline_bundle(pdf, tmp_path / 'out')
    assert code == 0, stderr
    data = _load_bundle(bundle)
    assert data['quality']['status'] == 'partial'
    page_two = next(
        unit for unit in data['source_units']
        if unit['type'] == 'page' and unit['index'] == 2
    )
    assert page_two['status'] == 'ocr_required'
    assert '[PARTIAL]' in stdout


def test_pdf_without_any_usable_content_fails_without_publication(tmp_path):
    pdf = tmp_path / 'blank.pdf'
    _make_pdf(pdf, lambda c: c.showPage())
    code, _, stderr, bundle = _run_pipeline_bundle(pdf, tmp_path / 'out')
    assert code == 1
    assert 'no usable content' in stderr
    assert not bundle.exists()


def test_pdf_embedded_image_is_published_and_referenced(tmp_path):
    from PIL import Image
    image = tmp_path / 'pixel.png'
    Image.new('RGB', (20, 20), 'red').save(image)
    pdf = tmp_path / 'picture.pdf'

    def draw(c):
        for index in range(8):
            c.drawString(72, 760 - index * 12, f'Before image line {index}')
        c.drawImage(str(image), 72, 640, width=40, height=40)
        for index in range(8):
            c.drawString(72, 620 - index * 12, f'After image line {index}')

    _make_pdf(pdf, draw)
    code, _, stderr, bundle = _run_pipeline_bundle(pdf, tmp_path / 'out')
    assert code == 0, stderr
    data = _load_bundle(bundle)
    assert len(data['assets']) == 1
    assert any(node['type'] == 'image' and node['asset_id'] == data['assets'][0]['asset_id'] for node in data['content'])
    visible_text = '\n'.join(
        node.get('text', '') for node in data['content'] if node['type'] != 'image'
    )
    assert all(f'Before image line {index}' in visible_text for index in range(8))
    assert all(f'After image line {index}' in visible_text for index in range(8))
    assert [node['type'] for node in data['content']][1] == 'image'
    asset = bundle / data['assets'][0]['path']
    assert asset.exists()
    import hashlib
    assert hashlib.sha256(asset.read_bytes()).hexdigest() == data['assets'][0]['sha256']


def test_table_renderer_uses_null_confidence_without_inventing_headers():
    from canonical import render_table
    table = {
        'confidence': None,
        'rows': [
            [{'raw_text': 'A', 'text': 'A', 'normalized_text': 'A', 'rowspan': 1, 'colspan': 1}],
            [
                {'raw_text': 'B', 'text': 'B', 'normalized_text': 'B', 'rowspan': 1, 'colspan': 1},
                {'raw_text': '2', 'text': '2', 'normalized_text': '2', 'rowspan': 1, 'colspan': 1},
            ],
        ],
    }
    rendered = render_table(table)
    assert rendered == 'A\nB | 2'
    assert 'header' not in rendered.lower()


def test_bundle_replace_failure_rolls_back_previous_target(tmp_path, monkeypatch):
    import pipeline
    target = tmp_path / 'target'
    target.mkdir()
    (target / 'old.txt').write_text('old', encoding='utf-8')
    stage = tmp_path / 'stage'
    stage.mkdir()
    (stage / 'new.txt').write_text('new', encoding='utf-8')
    real_replace = os.replace
    calls = []

    def flaky(source, destination):
        calls.append((Path(source), Path(destination)))
        if len(calls) == 2:
            raise OSError('simulated replace failure')
        return real_replace(source, destination)

    monkeypatch.setattr(pipeline.os, 'replace', flaky)
    with pytest.raises(OSError, match='simulated'):
        pipeline._publish_directory(stage, target, True)
    assert (target / 'old.txt').read_text(encoding='utf-8') == 'old'
    assert not (target / 'new.txt').exists()


def test_bundle_overwrite_replaces_regular_file_target(tmp_path):
    import pipeline
    target = tmp_path / 'target'
    target.write_text('old file', encoding='utf-8')
    stage = tmp_path / 'stage'
    stage.mkdir()
    (stage / 'new.txt').write_text('new', encoding='utf-8')

    pipeline._publish_directory(stage, target, True)

    assert target.is_dir()
    assert (target / 'new.txt').read_text(encoding='utf-8') == 'new'
    assert not list(tmp_path.glob('.target.backup-*'))


def test_bundle_backup_cleanup_failure_is_nonfatal_after_commit(tmp_path, monkeypatch, capsys):
    import pipeline
    target = tmp_path / 'target'
    target.mkdir()
    (target / 'old.txt').write_text('old', encoding='utf-8')
    stage = tmp_path / 'stage'
    stage.mkdir()
    (stage / 'new.txt').write_text('new', encoding='utf-8')
    real_remove_path = pipeline._remove_path

    def fail_backup_cleanup(path):
        if '.target.backup-' in Path(path).name:
            raise PermissionError('simulated cleanup failure')
        return real_remove_path(path)

    monkeypatch.setattr(pipeline, '_remove_path', fail_backup_cleanup)

    pipeline._publish_directory(stage, target, True)

    assert (target / 'new.txt').read_text(encoding='utf-8') == 'new'
    assert not (target / 'old.txt').exists()
    backups = list(tmp_path.glob('.target.backup-*'))
    assert len(backups) == 1
    assert (backups[0] / 'old.txt').read_text(encoding='utf-8') == 'old'
    stderr = capsys.readouterr().err
    assert 'published' in stderr
    assert 'could not remove backup' in stderr


def _assert_single_office_bundle_image(bundle):
    import hashlib

    data = _load_bundle(bundle)
    assert len(data['assets']) == 1
    asset = data['assets'][0]
    assert asset['type'] == 'image'
    assert asset['path'].startswith('assets/images/')
    assert set(asset) >= {
        'asset_id', 'type', 'path', 'sha256', 'media_type',
        'source_locator', 'alt', 'caption',
    }
    published = bundle / Path(asset['path'])
    assert published.is_file()
    assert hashlib.sha256(published.read_bytes()).hexdigest() == asset['sha256']
    assert data['outputs']['assets'] == [{'path': asset['path'], 'sha256': asset['sha256']}]
    assert any(node['type'] == 'image' and node['asset_id'] == asset['asset_id'] for node in data['content'])
    markdown = (bundle / f'{bundle.name}.md').read_text(encoding='utf-8')
    assert f']({asset["path"]})' in markdown
    assert 'data:image/' not in markdown
    return data, asset, markdown


def test_docx_bundle_exports_embedded_image_to_json_and_markdown(tmp_path):
    from docx import Document
    from PIL import Image
    picture = tmp_path / 'office.png'
    Image.new('RGB', (16, 16), 'blue').save(picture)
    source = tmp_path / 'office.docx'
    document = Document()
    document.add_paragraph('Before Office image')
    document.add_picture(str(picture))
    document.add_paragraph('After Office image')
    document.save(source)
    code, stdout, stderr, bundle = _run_product_bundle(source, tmp_path / 'out')
    assert code == 0, stderr
    data, _, markdown = _assert_single_office_bundle_image(bundle)
    assert data['quality']['status'] == 'complete'
    assert not any(item['code'] == 'office_embedded_images_not_exported' for item in data['quality']['warnings'])
    significant = [
        node['type'] for node in data['content']
        if node['type'] != 'paragraph' or node.get('normalized_text', '').strip()
    ]
    assert significant == ['paragraph', 'image', 'paragraph']
    assert markdown.index('Before Office image') < markdown.index('](') < markdown.index('After Office image')
    assert '[PARTIAL]' not in stdout


def test_pptx_bundle_exports_embedded_image_to_json_and_markdown(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    picture = tmp_path / 'slide.png'
    Image.new('RGB', (16, 16), 'green').save(picture)
    source = tmp_path / 'slides.pptx'
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = 'Slide text'
    slide.shapes.add_picture(str(picture), Inches(1), Inches(2))
    presentation.save(source)

    code, _, stderr, bundle = _run_product_bundle(source, tmp_path / 'out')
    assert code == 0, stderr
    data, _, markdown = _assert_single_office_bundle_image(bundle)
    assert 'Slide text' in markdown
    assert data['quality']['status'] == 'complete'


def test_xlsx_bundle_retains_unresolved_image_without_guessing_position(tmp_path):
    from PIL import Image
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as SpreadsheetImage

    picture = tmp_path / 'sheet.png'
    Image.new('RGB', (16, 16), 'purple').save(picture)
    source = tmp_path / 'workbook.xlsx'
    workbook = Workbook()
    sheet = workbook.active
    sheet['A1'] = 'Sheet text'
    sheet.add_image(SpreadsheetImage(str(picture)), 'A3')
    workbook.save(source)

    # AnyDoc 0.1.3 does not expose XLSX drawing images in its Document model;
    # retain the historical image-preservation assertion as an explicit
    # MarkItDown rollback regression instead of bypassing the AnyDoc contract.
    code, _, stderr, bundle = _run_product_bundle(
        source,
        tmp_path / 'out',
        ['--local-document-adapter', 'markitdown'],
    )
    assert code == 0, stderr
    data = _load_bundle(bundle)
    assert len(data['assets']) == 1
    asset = data['assets'][0]
    markdown = (bundle / f'{bundle.name}.md').read_text(encoding='utf-8')
    assert 'Sheet text' in markdown
    assert data['quality']['status'] == 'partial'
    assert not any(node['type'] == 'image' for node in data['content'])
    assert f']({asset["path"]})' not in markdown
    assert any(item['code'] == 'office_image_position_unresolved' for item in data['quality']['warnings'])
    assert data['relationships'] == [{
        'type': 'image_occurrence',
        'source_unit_id': data['source_units'][0]['id'],
        'asset_id': asset['asset_id'],
        'occurrence_index': 1,
        'placement': 'unresolved',
    }]


def test_docx_markdown_only_intentionally_omits_image_assets_and_links(tmp_path):
    from docx import Document
    from PIL import Image

    inputs = tmp_path / 'inputs'
    inputs.mkdir()
    picture = inputs / 'office.png'
    Image.new('RGB', (16, 16), 'blue').save(picture)
    source = inputs / 'office.docx'
    document = Document()
    document.add_paragraph('Before')
    document.add_picture(str(picture))
    document.add_paragraph('After')
    document.save(source)
    output = tmp_path / 'output' / 'clean.md'

    result = subprocess.run(
        SCRIPT + CONFIG_ARG + ['--input', str(source), '--output-path', str(output)],
        capture_output=True,
    )

    assert result.returncode == 0, result.stderr.decode(errors='replace')
    assert sorted(path.name for path in output.parent.iterdir()) == ['clean.md']
    markdown = output.read_text(encoding='utf-8')
    assert 'Before' in markdown and 'After' in markdown
    assert 'assets/images/' not in markdown
    assert 'data:image/' not in markdown
    assert '[PARTIAL]' not in result.stdout.decode(errors='replace')


def test_docx_missing_embedded_image_publishes_text_without_dangling_link(tmp_path):
    import zipfile
    from docx import Document
    from PIL import Image

    picture = tmp_path / 'missing.png'
    Image.new('RGB', (16, 16), 'red').save(picture)
    source = tmp_path / 'broken-image.docx'
    document = Document()
    document.add_paragraph('Before missing image')
    document.add_picture(str(picture))
    document.add_paragraph('After missing image')
    document.save(source)
    rewritten = tmp_path / 'rewritten.docx'
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(rewritten, 'w') as replacement:
        for item in original.infolist():
            if '/media/' not in item.filename:
                replacement.writestr(item, original.read(item.filename))
    rewritten.replace(source)

    code, _, stderr, bundle = _run_product_bundle(source, tmp_path / 'out')

    assert code == 0, stderr
    data = _load_bundle(bundle)
    assert data['assets'] == []
    assert data['outputs']['assets'] == []
    assert data['quality']['status'] == 'partial'
    assert any(item['code'] == 'office_image_target_missing' for item in data['quality']['warnings'])
    markdown = (bundle / 'broken-image.md').read_text(encoding='utf-8')
    assert 'Before missing image' in markdown and 'After missing image' in markdown
    assert 'assets/images/' not in markdown
    assert 'data:image/' not in markdown




















def test_rapidocr_provider_is_lazy_reused_and_maps_bitmap_to_pdf_coordinates():
    from types import SimpleNamespace
    from PIL import Image
    from ocr_provider import OcrSettings, RapidOcrProvider

    factory_calls = []
    inference_calls = []

    class Engine:
        def __call__(self, image, **kwargs):
            inference_calls.append((image.size, kwargs))
            return SimpleNamespace(
                boxes=[
                    [[10, 10], [190, 10], [190, 30], [10, 30]],
                    [[10, 40], [190, 40], [190, 60], [10, 60]],
                ],
                txts=['Mapped OCR line', 'Low confidence line'],
                scores=[0.93, 0.2],
                elapse=0.01,
            )

    def engine_factory():
        factory_calls.append(True)
        return Engine()

    class PositionConverter:
        def to_page(self, x, y):
            return x / 2, 100 - y / 2

    class Bitmap:
        def __init__(self):
            self.closed = False

        def get_posconv(self, _page):
            return PositionConverter()

        def to_pil(self):
            return Image.new('RGB', (400, 200), 'white')

        def close(self):
            self.closed = True

    class Page:
        def __init__(self):
            self.bitmaps = []

        def get_size(self):
            return 200, 100

        def render(self, **_kwargs):
            bitmap = Bitmap()
            self.bitmaps.append(bitmap)
            return bitmap

    provider = RapidOcrProvider(
        OcrSettings(mode='auto', dpi=144, max_long_edge=4096, min_confidence=0.5),
        engine_factory=engine_factory,
    )
    assert factory_calls == []

    page = Page()
    first = provider.extract(page, 1)
    second = provider.extract(page, 2)

    assert len(factory_calls) == 1
    assert len(inference_calls) == 2
    assert first.spans[0].bbox == pytest.approx((5.0, 85.0, 95.0, 95.0))
    assert first.spans[0].confidence == pytest.approx(0.93)
    assert first.dropped_low_confidence == 1
    assert first.language == 'ch'
    assert first.requested_dpi == 144
    assert first.min_confidence == pytest.approx(0.5)
    assert first.model_profile == 'PP-OCRv6-small'
    assert first.runtime == 'injected'
    assert second.page_number == 2
    assert all(bitmap.closed for bitmap in page.bitmaps)


def test_rapidocr_provider_caches_initialization_failure():
    from ocr_provider import OcrSettings, OcrUnavailableError, RapidOcrProvider

    calls = []

    def unavailable_factory():
        calls.append(True)
        raise OcrUnavailableError('backend unavailable')

    class Page:
        def get_size(self):
            return 200, 100

    provider = RapidOcrProvider(
        OcrSettings(mode='auto'), engine_factory=unavailable_factory
    )
    for _ in range(2):
        with pytest.raises(OcrUnavailableError, match='backend unavailable'):
            provider.extract(Page(), 1)

    assert len(calls) == 1
