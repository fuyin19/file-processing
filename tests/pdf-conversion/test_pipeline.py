from __future__ import annotations

import importlib.util
import json
import os
from pathlib import Path
import shutil
import subprocess
import sys
import zipfile

import pytest
from pypdf import PdfReader, PdfWriter


ROOT = Path(__file__).resolve().parents[2]
SCRIPT = ROOT / "skills" / "pdf-conversion" / "scripts" / "pipeline.py"
SHARED = ROOT / "skills" / "_shared" / "scripts"
MARKDOWN = ROOT / "skills" / "markdown-conversion" / "scripts"
CONFIG = ROOT / "tests" / "pdf-conversion" / "fixtures" / "test_config.json"
LO = Path(r"C:\Program Files\LibreOffice\program\soffice.com")

for value in (str(SHARED), str(MARKDOWN)):
    if value not in sys.path:
        sys.path.insert(0, value)


def _module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


pipeline = _module("pdf_conversion_pipeline_tests", SCRIPT)
import conversion_runtime
import libreoffice_pdf


def _cli(*args: object, timeout: float = 120) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(SCRIPT), "--config", str(CONFIG), *(str(item) for item in args)],
        cwd=ROOT,
        capture_output=True,
        text=True,
        encoding="utf-8",
        timeout=timeout,
    )


def _pdf(path: Path, pages: int = 1, *, encrypted: bool = False) -> Path:
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=612, height=792)
    if encrypted:
        writer.encrypt("secret")
    with path.open("wb") as stream:
        writer.write(stream)
    return path


def _corrupt_filtered_content_pdf(path: Path) -> Path:
    from pypdf.generic import NameObject, StreamObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=100, height=100)
    contents = StreamObject()
    contents._data = b"not-hex-data>"
    contents[NameObject("/Filter")] = NameObject("/ASCIIHexDecode")
    page[NameObject("/Contents")] = writer._add_object(contents)
    with path.open("wb") as stream:
        writer.write(stream)
    return path


def _fake_args(source: Path, output: Path, *extra: str):
    args = pipeline.build_parser().parse_args(
        ["--input", str(source), "--output-dir", str(output), *extra]
    )
    pipeline.precheck(args)
    return args


class FakeEngine:
    def __init__(self, valid_pdf: Path):
        self.valid_pdf = valid_pdf
        self.settings = libreoffice_pdf.PdfConversionSettings()
        self.calls = []

    def convert(self, snapshot, destination, workspace):
        self.calls.append((snapshot, Path(destination), Path(workspace)))
        Path(destination).parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(self.valid_pdf, destination)
        return {"sha256": "a" * 64, "size_bytes": Path(destination).stat().st_size, "pages": 1}


def test_frontmatter_and_versions_are_coherent():
    skill = (ROOT / "skills" / "pdf-conversion" / "SKILL.md").read_text(encoding="utf-8")
    assert "name: pdf-conversion" in skill
    assert "version: 2.0.0" in skill
    assert pipeline.VERSION == "2.0.0"


def test_exact_filter_argument_is_one_sorted_ascii_json_value():
    value = libreoffice_pdf.filter_argument("writer")
    prefix = "pdf:writer_pdf_Export:"
    assert value.startswith(prefix)
    payload = value[len(prefix):]
    assert payload == json.dumps(json.loads(payload), sort_keys=True, separators=(",", ":"), ensure_ascii=True)
    properties = json.loads(payload)
    assert properties["SelectPdfVersion"] == {"type": "long", "value": "17"}
    assert properties["UseTaggedPDF"] == {"type": "boolean", "value": "true"}
    assert properties["IsSkipEmptyPages"]["value"] == "false"


@pytest.mark.parametrize(
    ("family", "specific"),
    [
        ("writer", {"ExportNotesInMargin", "IsSkipEmptyPages"}),
        ("impress", {"ExportHiddenSlides", "ExportNotesPages", "ExportOnlyNotesPages", "UseTransitionEffects"}),
        ("calc", {"SinglePageSheets"}),
    ],
)
def test_filter_family_properties(family, specific):
    properties = libreoffice_pdf.filter_properties(family)
    assert specific <= properties.keys()
    assert properties["ReduceImageResolution"]["value"] == "false"
    assert properties["UseLosslessCompression"]["value"] == "true"
    assert properties["ExportFormFields"]["value"] == "false"


def test_discovery_authoritative_cli_precedes_config(monkeypatch, tmp_path):
    cli = tmp_path / "cli" / "soffice.com"
    config = tmp_path / "config" / "soffice.com"
    cli.parent.mkdir()
    config.parent.mkdir()
    cli.write_bytes(b"cli")
    config.write_bytes(b"config")
    monkeypatch.setattr(libreoffice_pdf, "_probe_candidate", lambda path: "26.2" if path == cli else None)
    assert libreoffice_pdf.resolve_libreoffice(str(cli), str(config)) == (cli, "26.2")


def test_authoritative_exe_normalizes_to_com(monkeypatch, tmp_path):
    com = tmp_path / "soffice.com"
    com.write_bytes(b"x")
    monkeypatch.setattr(libreoffice_pdf, "_probe_candidate", lambda path: "1.2.3" if path == com else None)
    assert libreoffice_pdf.resolve_libreoffice(str(tmp_path / "soffice.exe"), "")[0] == com


def test_invalid_authoritative_path_fails_instead_of_falling_back(tmp_path):
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="Invalid authoritative"):
        libreoffice_pdf.resolve_libreoffice(str(tmp_path / "missing"), "")


def test_profile_hardening_contains_required_controls(tmp_path):
    profile = tmp_path / "profile"
    profile.mkdir()
    libreoffice_pdf.seed_profile(profile)
    text = (profile / "user" / "registrymodifications.xcu").read_text(encoding="utf-8")
    assert "MacroSecurityLevel" in text and ">3<" in text
    assert "SecureURL" in text
    assert "AutoCheckEnabled" in text and "false" in text
    assert "Writer/Content/Update" in text and "Calc/Content/Update" in text


def test_engine_builds_exact_single_filter_argv(tmp_path, monkeypatch):
    source = tmp_path / "source.docx"
    source.write_bytes(b"office")
    snapshot = conversion_runtime.acquire_source_snapshot(source, tmp_path / "snapshot" / "source.docx")
    valid = _pdf(tmp_path / "valid.pdf")
    captured = []

    def runner(argv, **kwargs):
        captured.append((argv, kwargs))
        output = Path(argv[argv.index("--outdir") + 1])
        output.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(valid, output / "source.pdf")
        return libreoffice_pdf.ProcessResult(0, b"ok", b"")

    monkeypatch.setattr(libreoffice_pdf, "_run_job_process", runner)
    monkeypatch.setattr(libreoffice_pdf, "validate_pdf", lambda path, settings: {"sha256": "a" * 64, "pages": 1, "size_bytes": path.stat().st_size})
    engine = libreoffice_pdf.LibreOfficePdfEngine(libreoffice_pdf.PdfConversionSettings())
    engine._resolved = (LO, "26.2.5.2")
    destination = tmp_path / "result" / "result.pdf"
    engine.convert(snapshot, destination, tmp_path / "private")
    argv, kwargs = captured[0]
    assert argv[0] == str(LO)
    assert argv[1].startswith("-env:UserInstallation=file:")
    assert argv[2:7] == ["--headless", "--nologo", "--nodefault", "--norestore", "--convert-to"]
    assert argv[7] == libreoffice_pdf.filter_argument("writer")
    assert argv[8] == "--outdir"
    assert len(argv) == 11
    assert kwargs["cwd"].name == "work"
    assert kwargs["environment"]["TEMP"] == kwargs["environment"]["TMP"]


def test_job_timeout_is_killable_and_bounded(tmp_path):
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="exceeded"):
        libreoffice_pdf._run_job_process(
            [sys.executable, "-c", "import time; time.sleep(30)"],
            cwd=tmp_path,
            environment=os.environ,
            timeout=0.05,
        )


def test_output_path_implies_direct_pdf_and_requires_pdf_suffix(tmp_path):
    source = _pdf(tmp_path / "source.pdf")
    args = pipeline.build_parser().parse_args(["--input", str(source), "--output-path", str(tmp_path / "out.pdf")])
    pipeline.precheck(args)
    assert args.output_mode == "pdf"
    bad = _cli("--input", source, "--output-path", tmp_path / "out.txt")
    assert bad.returncode == 1 and "must end in .pdf" in bad.stderr


def test_explicit_bundle_conflicts_with_output_path(tmp_path):
    source = _pdf(tmp_path / "source.pdf")
    result = _cli("--input", source, "--output-mode", "bundle", "--output-path", tmp_path / "out.pdf")
    assert result.returncode == 1
    assert "conflicts" in result.stderr


@pytest.mark.parametrize("value", ["https://example.com/a.pdf", "ftp://example.com/a.docx"])
def test_urls_are_rejected(value):
    result = _cli("--input", value)
    assert result.returncode == 1 and "URLs are not supported" in result.stderr


@pytest.mark.parametrize("suffix", [".dotx", ".potm", ".xltx", ".ppsm", ".txt"])
def test_unsupported_and_template_suffixes_are_rejected(tmp_path, suffix):
    source = tmp_path / f"source{suffix}"
    source.write_bytes(b"x")
    result = _cli("--input", source)
    assert result.returncode == 1 and "supported" in result.stderr.lower()


def test_pdf_source_direct_alias_is_rejected_even_with_overwrite(tmp_path):
    source = _pdf(tmp_path / "source.pdf")
    result = _cli("--input", source, "--output-mode", "pdf", "--overwrite")
    assert result.returncode == 1
    assert "alias" in result.stderr.lower()
    assert PdfReader(source).pages


def test_pdf_input_bundle_is_exact_copy_plus_source(tmp_path):
    source = _pdf(tmp_path / "source.pdf", 2)
    out = tmp_path / "out"
    result = _cli("--input", source, "--output-dir", out)
    assert result.returncode == 0, result.stderr
    bundle = out / "source"
    assert (bundle / "source.pdf").read_bytes() == source.read_bytes()
    assert (bundle / "src" / "source.pdf").read_bytes() == source.read_bytes()


@pytest.mark.parametrize("stem", [".cortex", ".CoRtEx-item"])
def test_pdf_bundle_rejects_reserved_cortex_stem_before_write(tmp_path, stem):
    source = _pdf(tmp_path / f"{stem}.pdf")
    output = tmp_path / "out"

    result = _cli("--input", source, "--output-dir", output, "--overwrite")

    assert result.returncode == 1
    assert "reserved Cortex name" in result.stderr
    assert not output.exists()


def test_pdf_only_record_bundle_remains_supported(tmp_path):
    source = _pdf(tmp_path / "record.pdf")
    output = tmp_path / "out"

    result = _cli("--input", source, "--output-dir", output)

    assert result.returncode == 0, result.stderr
    assert (output / "record/record.pdf").read_bytes() == source.read_bytes()


def test_direct_mode_emits_exactly_one_pdf(tmp_path):
    source = _pdf(tmp_path / "input.pdf")
    result = _cli("--input", source, "--output-mode", "pdf", "--output-path", tmp_path / "named.pdf")
    assert result.returncode == 0, result.stderr
    assert [path.name for path in tmp_path.iterdir() if path.name != "input.pdf"] == ["named.pdf"]


def test_pdf_bundle_supports_long_windows_source_and_output_paths(tmp_path):
    short = _pdf(tmp_path / "short.pdf")
    long_name = f"{'x' * 220}.pdf"
    source = tmp_path / long_name
    libreoffice_pdf.np.copy_file(short, source)
    out = tmp_path / "out"
    result = _cli("--input", source, "--output-dir", out)
    assert result.returncode == 0, result.stderr
    bundle = out / source.stem
    assert libreoffice_pdf.np.is_file(bundle / long_name)
    assert libreoffice_pdf.np.is_file(bundle / "src" / long_name)


def test_collision_exit_two_rename_changes_stem_and_overwrite_replaces(tmp_path):
    source = _pdf(tmp_path / "source.pdf")
    out = tmp_path / "out"
    assert _cli("--input", source, "--output-dir", out).returncode == 0
    assert _cli("--input", source, "--output-dir", out).returncode == 2
    renamed = _cli("--input", source, "--output-dir", out, "--rename")
    assert renamed.returncode == 0
    assert (out / "source_1" / "source_1.pdf").is_file()
    (out / "source" / "sentinel").write_text("old", encoding="utf-8")
    replaced = _cli("--input", source, "--output-dir", out, "--overwrite")
    assert replaced.returncode == 0
    assert not (out / "source" / "sentinel").exists()


def test_office_collision_is_detected_before_libreoffice_discovery(tmp_path, monkeypatch):
    source = tmp_path / "source.docx"
    source.write_bytes(b"office")
    out = tmp_path / "out"
    (out / "source").mkdir(parents=True)
    args = _fake_args(source, out)
    monkeypatch.setattr(
        libreoffice_pdf,
        "resolve_libreoffice",
        lambda *a, **k: pytest.fail("LibreOffice discovery ran before collision preflight"),
    )
    with pytest.raises(conversion_runtime.OutputCollision):
        pipeline.convert_one(args, str(source), pipeline.load_config(CONFIG))


def test_batch_layout_types_and_collision_exit(tmp_path):
    inputs = tmp_path / "inputs"
    (inputs / "nested").mkdir(parents=True)
    _pdf(inputs / "a.pdf")
    _pdf(inputs / "nested" / "b.pdf")
    (inputs / "ignored.txt").write_text("x", encoding="utf-8")
    out = tmp_path / "out"
    result = _cli("--input-dir", inputs, "--output-dir", out, "--types", "pdf")
    assert result.returncode == 0, result.stderr
    assert (out / "a" / "a.pdf").is_file()
    assert (out / "nested" / "b" / "b.pdf").is_file()
    assert _cli("--input-dir", inputs, "--output-dir", out, "--types", "pdf").returncode == 2


def test_snapshot_write_injection_is_denied_and_bytes_stable(tmp_path):
    source = tmp_path / "source.docx"
    source.write_bytes(b"original")
    destination = tmp_path / "stage" / "source.docx"

    def attempt_write(src, dst):
        with pytest.raises((PermissionError, OSError)):
            Path(src).write_bytes(b"mutated")
        shutil.copyfile(src, dst)

    snapshot = conversion_runtime.acquire_source_snapshot(source, destination, compatibility_copier=attempt_write)
    assert source.read_bytes() == destination.read_bytes() == b"original"
    snapshot.verify()


def test_snapshot_replacement_injection_is_denied_and_identity_stable(tmp_path):
    source = tmp_path / "source.docx"
    source.write_bytes(b"original")
    replacement = tmp_path / "replacement.docx"
    replacement.write_bytes(b"replacement")

    def attempt_replace(src, dst):
        with pytest.raises((PermissionError, OSError)):
            os.replace(replacement, src)
        shutil.copyfile(src, dst)

    snapshot = conversion_runtime.acquire_source_snapshot(
        source,
        tmp_path / "stage" / "source.docx",
        compatibility_copier=attempt_replace,
    )
    assert source.read_bytes() == snapshot.physical_path.read_bytes() == b"original"
    snapshot.verify()


def test_pdf_validator_uses_one_locked_identity_pinned_stream(tmp_path, monkeypatch):
    import pypdf

    worker = _module(
        "pdf_validation_worker_lock_tests",
        SHARED / "pdf_validation_worker.py",
    )
    source = _pdf(tmp_path / "source.pdf")
    replacement = _pdf(tmp_path / "replacement.pdf", 2)
    real_reader = pypdf.PdfReader

    def injecting_reader(stream, *args, **kwargs):
        with pytest.raises((PermissionError, OSError)):
            source.write_bytes(b"mutated")
        with pytest.raises((PermissionError, OSError)):
            os.replace(replacement, source)
        return real_reader(stream, *args, **kwargs)

    monkeypatch.setattr(pypdf, "PdfReader", injecting_reader)
    result = worker.validate(source, 1024 * 1024)
    assert result["pages"] == 1
    assert PdfReader(source).pages


def test_shared_overwrite_commit_failure_restores_exact_old_target(tmp_path, monkeypatch):
    target = tmp_path / "target"
    target.mkdir()
    (target / "old.txt").write_text("old", encoding="utf-8")
    stage = conversion_runtime.new_owned_dir(tmp_path, ".test-stage-")
    (stage.path / "new.txt").write_text("new", encoding="utf-8")
    real = conversion_runtime.np.rename_no_replace

    def fail_commit(source, destination):
        if conversion_runtime.np.paths_equal(source, stage.path) and conversion_runtime.np.paths_equal(destination, target):
            raise OSError("injected commit failure")
        return real(source, destination)

    monkeypatch.setattr(conversion_runtime.np, "rename_no_replace", fail_commit)
    with pytest.raises(OSError, match="injected commit failure"):
        conversion_runtime.publish_owned(stage, target, True)
    assert (target / "old.txt").read_text(encoding="utf-8") == "old"
    assert not stage.path.exists()
    assert not list(tmp_path.glob(".conversion-backup-*"))


def test_pdf_postcommit_mutation_rolls_back_and_restores_old_target(tmp_path, monkeypatch):
    target = tmp_path / "target"
    target.mkdir()
    (target / "old.txt").write_text("old", encoding="utf-8")
    stage = conversion_runtime.new_owned_dir(tmp_path, ".test-stage-")
    staged_pdf = _pdf(stage.path / "document.pdf")
    settings = libreoffice_pdf.ValidationSettings()
    expected = libreoffice_pdf.validate_pdf(staged_pdf, settings)
    real_rename = conversion_runtime.np.rename_no_replace

    def inject_after_commit(source, destination):
        real_rename(source, destination)
        if conversion_runtime.np.paths_equal(source, stage.path) and conversion_runtime.np.paths_equal(destination, target):
            (target / "document.pdf").write_bytes(b"%PDF-1.7\ninvalid\n%%EOF\n")

    monkeypatch.setattr(conversion_runtime.np, "rename_no_replace", inject_after_commit)
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="PDF"):
        conversion_runtime.publish_owned(
            stage,
            target,
            True,
            verify_payload=lambda root: libreoffice_pdf.verify_validated_pdf(
                root / "document.pdf", settings, expected
            ),
        )
    assert (target / "old.txt").read_text(encoding="utf-8") == "old"
    assert not (target / "document.pdf").exists()
    assert not stage.path.exists()
    assert not list(tmp_path.glob(".conversion-backup-*"))


def test_ambiguous_fresh_commit_mutation_is_verified_and_rolled_back(tmp_path, monkeypatch):
    target = tmp_path / "target"
    stage = conversion_runtime.new_owned_dir(tmp_path, ".test-stage-")
    staged_pdf = _pdf(stage.path / "document.pdf")
    settings = libreoffice_pdf.ValidationSettings()
    expected = libreoffice_pdf.validate_pdf(staged_pdf, settings)
    real_rename = conversion_runtime.np.rename_no_replace

    def move_mutate_and_raise(source, destination):
        real_rename(source, destination)
        if conversion_runtime.np.paths_equal(source, stage.path) and conversion_runtime.np.paths_equal(destination, target):
            published_pdf = target / "document.pdf"
            inode = published_pdf.stat().st_ino
            published_pdf.write_bytes(b"%PDF-1.7\ninvalid\n%%EOF\n")
            assert published_pdf.stat().st_ino == inode
            raise OSError("injected ambiguous commit")

    monkeypatch.setattr(conversion_runtime.np, "rename_no_replace", move_mutate_and_raise)
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="PDF"):
        conversion_runtime.publish_owned(
            stage,
            target,
            False,
            verify_payload=lambda root: libreoffice_pdf.verify_validated_pdf(
                root / "document.pdf", settings, expected
            ),
        )
    assert not target.exists()
    assert not stage.path.exists()


def test_ambiguous_fresh_fileexists_mutation_is_verified_and_rolled_back(tmp_path, monkeypatch):
    target = tmp_path / "target"
    stage = conversion_runtime.new_owned_dir(tmp_path, ".test-stage-")
    staged_pdf = _pdf(stage.path / "document.pdf")
    settings = libreoffice_pdf.ValidationSettings()
    expected = libreoffice_pdf.validate_pdf(staged_pdf, settings)
    real_rename = conversion_runtime.np.rename_no_replace

    def move_mutate_and_raise_fileexists(source, destination):
        real_rename(source, destination)
        if conversion_runtime.np.paths_equal(source, stage.path) and conversion_runtime.np.paths_equal(destination, target):
            published_pdf = target / "document.pdf"
            inode = published_pdf.stat().st_ino
            published_pdf.write_bytes(b"%PDF-1.7\ninvalid\n%%EOF\n")
            assert published_pdf.stat().st_ino == inode
            raise FileExistsError("injected ambiguous FileExistsError")

    monkeypatch.setattr(
        conversion_runtime.np,
        "rename_no_replace",
        move_mutate_and_raise_fileexists,
    )
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="PDF"):
        conversion_runtime.publish_owned(
            stage,
            target,
            False,
            verify_payload=lambda root: libreoffice_pdf.verify_validated_pdf(
                root / "document.pdf", settings, expected
            ),
        )
    assert not target.exists()
    assert not stage.path.exists()
    assert conversion_runtime._owned_key(stage.path) not in conversion_runtime._OWNED_ENTRIES


def test_ambiguous_overwrite_commit_mutation_restores_old_target(tmp_path, monkeypatch):
    target = tmp_path / "target"
    target.mkdir()
    (target / "old.txt").write_text("old", encoding="utf-8")
    stage = conversion_runtime.new_owned_dir(tmp_path, ".test-stage-")
    staged_pdf = _pdf(stage.path / "document.pdf")
    settings = libreoffice_pdf.ValidationSettings()
    expected = libreoffice_pdf.validate_pdf(staged_pdf, settings)
    real_rename = conversion_runtime.np.rename_no_replace

    def move_mutate_and_raise(source, destination):
        real_rename(source, destination)
        if conversion_runtime.np.paths_equal(source, stage.path) and conversion_runtime.np.paths_equal(destination, target):
            published_pdf = target / "document.pdf"
            inode = published_pdf.stat().st_ino
            published_pdf.write_bytes(b"%PDF-1.7\ninvalid\n%%EOF\n")
            assert published_pdf.stat().st_ino == inode
            raise OSError("injected ambiguous commit")

    monkeypatch.setattr(conversion_runtime.np, "rename_no_replace", move_mutate_and_raise)
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="PDF"):
        conversion_runtime.publish_owned(
            stage,
            target,
            True,
            verify_payload=lambda root: libreoffice_pdf.verify_validated_pdf(
                root / "document.pdf", settings, expected
            ),
        )
    assert (target / "old.txt").read_text(encoding="utf-8") == "old"
    assert not (target / "document.pdf").exists()
    assert not stage.path.exists()
    assert not list(tmp_path.glob(".conversion-backup-*"))


def test_validator_rejects_corrupt_filtered_page_content(tmp_path):
    source = _corrupt_filtered_content_pdf(tmp_path / "corrupt.pdf")
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="content stream is unreadable"):
        libreoffice_pdf.validate_pdf(source, libreoffice_pdf.ValidationSettings())


def test_corrupt_filtered_pdf_input_publishes_nothing_and_preserves_overwrite(tmp_path):
    source = _corrupt_filtered_content_pdf(tmp_path / "corrupt.pdf")
    out = tmp_path / "out"
    target = out / "corrupt"
    target.mkdir(parents=True)
    (target / "old.txt").write_text("old", encoding="utf-8")

    result = _cli("--input", source, "--output-dir", out, "--overwrite")

    assert result.returncode == 1
    assert "content stream is unreadable" in result.stderr
    assert (target / "old.txt").read_text(encoding="utf-8") == "old"
    assert not (target / "corrupt.pdf").exists()


def test_malformed_encrypted_zero_page_and_over_limit_publish_nothing(tmp_path):
    sources = []
    malformed = tmp_path / "malformed.pdf"
    malformed.write_bytes(b"%PDF-1.7\nnot a pdf\n%%EOF\n")
    sources.append(malformed)
    sources.append(_pdf(tmp_path / "encrypted.pdf", encrypted=True))
    sources.append(_pdf(tmp_path / "zero.pdf", 0))
    out = tmp_path / "out"
    for source in sources:
        result = _cli("--input", source, "--output-dir", out)
        assert result.returncode == 1
        assert not (out / source.stem).exists()

    config = tmp_path / "small.json"
    config.write_text(json.dumps({"pdf_conversion": {"validation": {"max_pdf_bytes": 10}}}), encoding="utf-8")
    valid = _pdf(tmp_path / "large.pdf")
    result = subprocess.run(
        [sys.executable, str(SCRIPT), "--config", str(config), "--input", str(valid), "--output-dir", str(out)],
        cwd=ROOT, capture_output=True, text=True, encoding="utf-8",
    )
    assert result.returncode == 1 and "byte limit" in result.stderr
    assert not (out / "large").exists()


def test_fake_provider_failure_cleans_stage_and_preserves_old_target(tmp_path, monkeypatch):
    source = tmp_path / "source.docx"
    source.write_bytes(b"source")
    out = tmp_path / "out"
    target = out / "source"
    target.mkdir(parents=True)
    (target / "old").write_text("old", encoding="utf-8")
    args = _fake_args(source, out, "--overwrite")
    monkeypatch.setattr(FakeEngine, "convert", lambda *a, **k: (_ for _ in ()).throw(libreoffice_pdf.LibreOfficeError("timeout")))
    with pytest.raises(libreoffice_pdf.LibreOfficeError, match="timeout"):
        pipeline.convert_one(args, str(source), pipeline.load_config(CONFIG), engine=FakeEngine(_pdf(tmp_path / "valid.pdf")))
    assert (target / "old").read_text(encoding="utf-8") == "old"
    assert not list(out.glob(".pc-stage-*"))


def _copy_as(source: Path, destination: Path) -> Path:
    destination.write_bytes(source.read_bytes())
    return destination


@pytest.mark.parametrize(
    ("suffix", "fixture"),
    [
        (".doc", "text.doc"), (".docx", "text.docx"), (".docm", "text.docx"),
        (".ppt", "pres.ppt"), (".pptx", "pres.pptx"), (".pptm", "pres.pptx"),
        (".pps", "pres.ppt"), (".ppsx", "pres.pptx"),
        (".xls", "sheet.xls"), (".xlsx", "sheet.xlsx"), (".xlsm", "sheet.xlsx"),
        (".xlsb", "any_sheets.xlsb"),
    ],
)
def test_real_libreoffice_supports_every_office_suffix(tmp_path, suffix, fixture):
    assert LO.is_file(), "Required installed LibreOffice is missing"
    source_fixture = ROOT / "tests" / "markdown-conversion" / "fixtures" / "anydoc" / fixture
    source = _copy_as(source_fixture, tmp_path / f"source{suffix}")
    result = _cli("--input", source, "--output-mode", "pdf", "--output-path", tmp_path / "out.pdf", timeout=180)
    assert result.returncode == 0, result.stderr
    assert len(PdfReader(tmp_path / "out.pdf").pages) >= 1


def test_real_writer_two_pages_text_bookmark_tagging_and_no_downsample(tmp_path):
    from docx import Document
    from docx.shared import Inches
    from PIL import Image

    image = tmp_path / "large.png"
    Image.new("RGB", (900, 700), (10, 90, 180)).save(image)
    document = Document()
    document.add_heading("Export Bookmark", level=1)
    document.add_paragraph("Writer first page sentinel")
    document.add_picture(str(image), width=Inches(5))
    document.add_page_break()
    document.add_paragraph("Writer second page sentinel")
    source = tmp_path / "writer.docx"
    document.save(source)
    result = _cli("--input", source, "--output-mode", "pdf", "--output-path", tmp_path / "writer.pdf", timeout=180)
    assert result.returncode == 0, result.stderr
    raw = (tmp_path / "writer.pdf").read_bytes()
    assert raw.startswith(b"%PDF-1.7")
    reader = PdfReader(tmp_path / "writer.pdf")
    assert len(reader.pages) == 2
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert "first page sentinel" in text and "second page sentinel" in text
    assert reader.outline
    root = reader.trailer["/Root"]
    assert bool(root.get("/MarkInfo", {}).get("/Marked"))
    assert not reader.get_fields()
    widths = []
    for page in reader.pages:
        resources = page.get("/Resources", {})
        xobjects = resources.get("/XObject", {})
        for ref in xobjects.values():
            obj = ref.get_object()
            if obj.get("/Subtype") == "/Image":
                widths.append(int(obj.get("/Width", 0)))
    assert max(widths) >= 900


def test_real_impress_excludes_hidden_slide_and_notes(tmp_path):
    from pptx import Presentation

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Visible slide sentinel"
    first.notes_slide.notes_text_frame.text = "Notes sentinel must not export"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Hidden slide sentinel"
    source = tmp_path / "slides.pptx"
    presentation.save(source)
    rewritten = tmp_path / "slides-hidden.pptx"
    with zipfile.ZipFile(source) as incoming, zipfile.ZipFile(rewritten, "w") as outgoing:
        for item in incoming.infolist():
            data = incoming.read(item.filename)
            if item.filename == "ppt/slides/slide2.xml":
                text = data.decode("utf-8")
                text = text.replace("<p:sld ", "<p:sld show=\"0\" ", 1)
                data = text.encode("utf-8")
            outgoing.writestr(item, data)
    result = _cli("--input", rewritten, "--output-mode", "pdf", "--output-path", tmp_path / "slides.pdf", timeout=180)
    assert result.returncode == 0, result.stderr
    reader = PdfReader(tmp_path / "slides.pdf")
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert len(reader.pages) == 1
    assert "Visible slide sentinel" in text
    assert "Hidden slide sentinel" not in text
    assert "Notes sentinel" not in text


def test_real_calc_honors_print_pages_and_hidden_sheet(tmp_path):
    from openpyxl import Workbook
    from openpyxl.worksheet.pagebreak import Break

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Visible"
    for row in range(1, 121):
        sheet.cell(row, 1, f"Visible row {row}")
    sheet.print_area = "A1:A120"
    sheet.row_breaks.append(Break(id=55))
    hidden = workbook.create_sheet("Hidden")
    hidden["A1"] = "Hidden sheet sentinel"
    hidden.sheet_state = "hidden"
    source = tmp_path / "book.xlsx"
    workbook.save(source)
    result = _cli("--input", source, "--output-mode", "pdf", "--output-path", tmp_path / "book.pdf", timeout=180)
    assert result.returncode == 0, result.stderr
    reader = PdfReader(tmp_path / "book.pdf")
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert len(reader.pages) >= 2
    assert "Visible row 1" in text and "Visible row 120" in text
    assert "Hidden sheet sentinel" not in text


def test_real_macro_named_input_does_not_mutate_sentinel(tmp_path):
    sentinel = tmp_path / "macro-sentinel.txt"
    sentinel.write_text("unchanged", encoding="utf-8")
    fixture = ROOT / "tests" / "markdown-conversion" / "fixtures" / "anydoc" / "text.docx"
    source = _copy_as(fixture, tmp_path / "macro.docm")
    result = _cli("--input", source, "--output-mode", "pdf", "--output-path", tmp_path / "macro.pdf", timeout=180)
    assert result.returncode == 0, result.stderr
    assert sentinel.read_text(encoding="utf-8") == "unchanged"
